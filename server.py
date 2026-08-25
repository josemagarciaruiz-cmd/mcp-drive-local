"""
MCP Drive — Servidor MCP con control TOTAL de Google Drive.

Cubre lo que el conector oficial de Claude NO hace: mover, renombrar,
enviar a papelera, restaurar, compartir (permisos), copiar, crear carpetas,
subir y descargar, y (con confirmación explícita) borrado definitivo.

Transporte: streamable-http (pensado para desplegar en un VPS, p.ej. Hostinger).
Autenticación: OAuth de usuario mediante refresh token (un solo usuario, el titular).

Variables de entorno necesarias (ver .env.example):
    GOOGLE_CLIENT_ID
    GOOGLE_CLIENT_SECRET
    GOOGLE_REFRESH_TOKEN
    MCP_HOST      (opcional, por defecto 0.0.0.0)
    MCP_PORT      (opcional, por defecto 8000)

Autor: José María García Ruiz · josemaria.ai
"""

import os
import base64
import io
from typing import Any, Optional

from mcp.server.fastmcp import FastMCP

from google.oauth2.credentials import Credentials
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
from googleapiclient.http import MediaIoBaseUpload, MediaIoBaseDownload

# --------------------------------------------------------------------------- #
# Configuración
# --------------------------------------------------------------------------- #

SCOPES = ["https://www.googleapis.com/auth/drive",
          "https://www.googleapis.com/auth/gmail.modify",
          "https://www.googleapis.com/auth/gmail.send",
          "https://www.googleapis.com/auth/gmail.settings.basic",
          "https://www.googleapis.com/auth/calendar",
          "https://www.googleapis.com/auth/tasks",
          "https://www.googleapis.com/auth/contacts",
          "https://www.googleapis.com/auth/forms.body",
          "https://www.googleapis.com/auth/forms.responses.readonly"]
TOKEN_URI = "https://oauth2.googleapis.com/token"

# Campos estándar que devolvemos de cada archivo/carpeta
FILE_FIELDS = (
    "id, name, mimeType, parents, size, modifiedTime, createdTime, "
    "trashed, webViewLink, owners(emailAddress)"
)

MCP_HOST = os.environ.get("MCP_HOST", "0.0.0.0")
MCP_PORT = int(os.environ.get("MCP_PORT", "8000"))

mcp = FastMCP("MCP Drive", host=MCP_HOST, port=MCP_PORT)

_service = None  # cache del cliente de la API
_docs_service = None  # cache del cliente de Docs


def _build_creds():
    """Credenciales OAuth compartidas (Drive y Docs; el scope drive cubre ambas)."""
    for n in ("GOOGLE_CLIENT_ID", "GOOGLE_CLIENT_SECRET", "GOOGLE_REFRESH_TOKEN"):
        if not os.environ.get(n):
            raise RuntimeError("Falta la variable de entorno " + n)
    return Credentials(
        token=None,
        refresh_token=os.environ["GOOGLE_REFRESH_TOKEN"],
        token_uri=TOKEN_URI,
        client_id=os.environ["GOOGLE_CLIENT_ID"],
        client_secret=os.environ["GOOGLE_CLIENT_SECRET"],
        scopes=SCOPES,
    )


def _get_docs():
    """Cliente de la API de Google Docs (edicion del cuerpo de documentos)."""
    global _docs_service
    if _docs_service is None:
        _docs_service = build("docs", "v1", credentials=_build_creds(), cache_discovery=False)
    return _docs_service


def _get_service():
    """Construye (y cachea) el cliente de la Drive API a partir del refresh token."""
    global _service
    if _service is not None:
        return _service

    client_id = os.environ.get("GOOGLE_CLIENT_ID")
    client_secret = os.environ.get("GOOGLE_CLIENT_SECRET")
    refresh_token = os.environ.get("GOOGLE_REFRESH_TOKEN")

    missing = [
        n for n, v in (
            ("GOOGLE_CLIENT_ID", client_id),
            ("GOOGLE_CLIENT_SECRET", client_secret),
            ("GOOGLE_REFRESH_TOKEN", refresh_token),
        ) if not v
    ]
    if missing:
        raise RuntimeError(
            "Faltan variables de entorno: " + ", ".join(missing) +
            ". Configúralas antes de arrancar el servidor (ver .env.example)."
        )

    creds = Credentials(
        token=None,
        refresh_token=refresh_token,
        token_uri=TOKEN_URI,
        client_id=client_id,
        client_secret=client_secret,
        scopes=SCOPES,
    )
    _service = build("drive", "v3", credentials=creds, cache_discovery=False)
    return _service


def _err(action: str, e: Exception) -> dict:
    """Formatea un error de forma accionable para el agente."""
    if isinstance(e, HttpError):
        try:
            reason = e.error_details
        except Exception:
            reason = str(e)
        return {
            "ok": False,
            "action": action,
            "error": f"La API de Drive devolvió {e.resp.status}: {reason}",
            "sugerencia": "Revisa el fileId/permiso. 404 = no existe o sin acceso; "
                          "403 = falta de permiso o cuota; 401 = token caducado.",
        }
    return {"ok": False, "action": action, "error": str(e)}


def _file(meta: dict) -> dict:
    """Normaliza la metadata de un archivo a un dict limpio."""
    return {
        "id": meta.get("id"),
        "name": meta.get("name"),
        "mimeType": meta.get("mimeType"),
        "parents": meta.get("parents", []),
        "size": meta.get("size"),
        "modifiedTime": meta.get("modifiedTime"),
        "trashed": meta.get("trashed"),
        "link": meta.get("webViewLink"),
    }


# --------------------------------------------------------------------------- #
# AUDITORIA (registro de acciones que modifican Drive)
# --------------------------------------------------------------------------- #

import json as _json
import time as _time
import threading as _threading

AUDIT_PATH = os.environ.get("AUDIT_LOG_PATH") or os.path.join(
    os.path.dirname(os.path.abspath(__file__)), "audit.jsonl")
_audit_lock = _threading.Lock()


def _audit(action: str, detail: dict) -> None:
    """Anexa una entrada JSONL al registro de auditoria. Nunca lanza."""
    try:
        rec = {"ts": _time.strftime("%Y-%m-%dT%H:%M:%SZ", _time.gmtime()),
               "action": action}
        rec.update(detail or {})
        d = os.path.dirname(AUDIT_PATH)
        if d:
            os.makedirs(d, exist_ok=True)
        with _audit_lock:
            with open(AUDIT_PATH, "a", encoding="utf-8") as f:
                f.write(_json.dumps(rec, ensure_ascii=False) + "\n")
    except Exception:
        pass


# --------------------------------------------------------------------------- #
# LECTURA
# --------------------------------------------------------------------------- #

@mcp.tool()
def drive_search(query: str, page_size: int = 50, page_token: Optional[str] = None) -> dict:
    """Busca archivos y carpetas en Drive con la sintaxis de consulta de la API.

    Ejemplos de `query`:
      - "name contains 'nomina'"
      - "'FOLDER_ID' in parents and trashed = false"
      - "mimeType = 'application/vnd.google-apps.folder'"
      - "fullText contains 'despido'"
    Devuelve hasta `page_size` resultados (máx. 1000).
    """
    try:
        service = _get_service()
        res = service.files().list(
            q=query,
            pageSize=max(1, min(page_size, 1000)),
            pageToken=page_token,
            fields=f"nextPageToken, files({FILE_FIELDS})",
            supportsAllDrives=True,
            includeItemsFromAllDrives=True,
        ).execute()
        files = [_file(f) for f in res.get("files", [])]
        return {"ok": True, "count": len(files), "files": files,
                "next_page_token": res.get("nextPageToken")}
    except Exception as e:
        return _err("drive_search", e)


@mcp.tool()
def drive_list_children(folder_id: str = "root", include_trashed: bool = False, page_token: Optional[str] = None) -> dict:
    """Lista el contenido directo de una carpeta. Usa 'root' para Mi unidad."""
    try:
        service = _get_service()
        q = f"'{folder_id}' in parents"
        if not include_trashed:
            q += " and trashed = false"
        res = service.files().list(
            q=q, pageSize=1000, pageToken=page_token,
            fields=f"nextPageToken, files({FILE_FIELDS})",
            supportsAllDrives=True, includeItemsFromAllDrives=True,
        ).execute()
        files = [_file(f) for f in res.get("files", [])]
        return {"ok": True, "folder_id": folder_id, "count": len(files), "files": files,
                "next_page_token": res.get("nextPageToken")}
    except Exception as e:
        return _err("drive_list_children", e)


@mcp.tool()
def drive_get_metadata(file_id: str) -> dict:
    """Devuelve la metadata de un archivo o carpeta por su id."""
    try:
        service = _get_service()
        meta = service.files().get(
            fileId=file_id, fields=FILE_FIELDS, supportsAllDrives=True
        ).execute()
        return {"ok": True, "file": _file(meta)}
    except Exception as e:
        return _err("drive_get_metadata", e)


# --------------------------------------------------------------------------- #
# ORGANIZACIÓN (lo que el conector oficial NO hace)
# --------------------------------------------------------------------------- #

@mcp.tool()
def drive_create_folder(name: str, parent_id: str = "root") -> dict:
    """Crea una carpeta nueva dentro de `parent_id` (por defecto Mi unidad)."""
    try:
        service = _get_service()
        body = {
            "name": name,
            "mimeType": "application/vnd.google-apps.folder",
            "parents": [parent_id],
        }
        meta = service.files().create(
            body=body, fields=FILE_FIELDS, supportsAllDrives=True
        ).execute()
        return {"ok": True, "file": _file(meta)}
    except Exception as e:
        return _err("drive_create_folder", e)


@mcp.tool()
def drive_move(file_id: str, new_parent_id: str, dry_run: bool = False) -> dict:
    """Mueve un archivo o carpeta a `new_parent_id` (quita el padre anterior).
    Con dry_run=True informa de lo que haria sin ejecutarlo."""
    try:
        service = _get_service()
        current = service.files().get(
            fileId=file_id, fields="parents, name", supportsAllDrives=True
        ).execute()
        prev_parents = ",".join(current.get("parents", []))
        if dry_run:
            return {"ok": True, "dry_run": True,
                    "haria": "mover '%s' a %s" % (current.get("name"), new_parent_id),
                    "desde": current.get("parents", []), "hacia": new_parent_id}
        meta = service.files().update(
            fileId=file_id,
            addParents=new_parent_id,
            removeParents=prev_parents,
            fields=FILE_FIELDS,
            supportsAllDrives=True,
        ).execute()
        _audit("drive_move", {"file_id": file_id, "name": meta.get("name"),
                              "desde": prev_parents, "hacia": new_parent_id})
        return {"ok": True, "file": _file(meta)}
    except Exception as e:
        return _err("drive_move", e)


@mcp.tool()
def drive_rename(file_id: str, new_name: str, dry_run: bool = False) -> dict:
    """Renombra un archivo o carpeta. Con dry_run=True no ejecuta, solo informa."""
    try:
        service = _get_service()
        if dry_run:
            cur = service.files().get(fileId=file_id, fields="name",
                                      supportsAllDrives=True).execute()
            return {"ok": True, "dry_run": True,
                    "haria": "renombrar '%s' -> '%s'" % (cur.get("name"), new_name)}
        meta = service.files().update(
            fileId=file_id, body={"name": new_name},
            fields=FILE_FIELDS, supportsAllDrives=True,
        ).execute()
        _audit("drive_rename", {"file_id": file_id, "new_name": new_name})
        return {"ok": True, "file": _file(meta)}
    except Exception as e:
        return _err("drive_rename", e)


@mcp.tool()
def drive_copy(file_id: str, new_name: Optional[str] = None,
               parent_id: Optional[str] = None) -> dict:
    """Copia un archivo. Opcionalmente con nombre nuevo y/o carpeta destino."""
    try:
        service = _get_service()
        body: dict[str, Any] = {}
        if new_name:
            body["name"] = new_name
        if parent_id:
            body["parents"] = [parent_id]
        meta = service.files().copy(
            fileId=file_id, body=body, fields=FILE_FIELDS, supportsAllDrives=True
        ).execute()
        return {"ok": True, "file": _file(meta)}
    except Exception as e:
        return _err("drive_copy", e)


# --------------------------------------------------------------------------- #
# PAPELERA Y BORRADO
# --------------------------------------------------------------------------- #

@mcp.tool()
def drive_trash(file_id: str, dry_run: bool = False) -> dict:
    """Envia un archivo o carpeta a la PAPELERA (reversible 30 dias).
    Con dry_run=True no ejecuta, solo informa."""
    try:
        service = _get_service()
        if dry_run:
            cur = service.files().get(fileId=file_id, fields="name",
                                      supportsAllDrives=True).execute()
            return {"ok": True, "dry_run": True,
                    "haria": "enviar a papelera '%s'" % cur.get("name")}
        meta = service.files().update(
            fileId=file_id, body={"trashed": True},
            fields=FILE_FIELDS, supportsAllDrives=True,
        ).execute()
        _audit("drive_trash", {"file_id": file_id, "name": meta.get("name")})
        return {"ok": True, "file": _file(meta)}
    except Exception as e:
        return _err("drive_trash", e)


@mcp.tool()
def drive_untrash(file_id: str) -> dict:
    """Restaura un archivo desde la papelera."""
    try:
        service = _get_service()
        meta = service.files().update(
            fileId=file_id, body={"trashed": False},
            fields=FILE_FIELDS, supportsAllDrives=True,
        ).execute()
        return {"ok": True, "file": _file(meta)}
    except Exception as e:
        return _err("drive_untrash", e)


@mcp.tool()
def drive_delete_permanent(file_id: str, confirm: bool = False) -> dict:
    """Borra DEFINITIVAMENTE un archivo (NO reversible). Requiere confirm=True.

    Úsalo solo con confirmación explícita del usuario; en la mayoría de casos
    es preferible drive_trash.
    """
    if not confirm:
        return {
            "ok": False,
            "action": "drive_delete_permanent",
            "error": "Borrado definitivo bloqueado. Vuelve a llamar con confirm=true "
                     "solo si el usuario lo ha pedido expresamente. Considera drive_trash.",
        }
    try:
        service = _get_service()
        service.files().delete(fileId=file_id, supportsAllDrives=True).execute()
        _audit("drive_delete_permanent", {"file_id": file_id})
        return {"ok": True, "deleted": file_id}
    except Exception as e:
        return _err("drive_delete_permanent", e)


# --------------------------------------------------------------------------- #
# COMPARTIR / PERMISOS
# --------------------------------------------------------------------------- #

@mcp.tool()
def drive_share(file_id: str, role: str = "reader", type: str = "user",
                email: Optional[str] = None, notify: bool = False) -> dict:
    """Crea un permiso sobre un archivo/carpeta.

    role: 'reader' | 'commenter' | 'writer' | 'owner'
    type: 'user' | 'group' | 'domain' | 'anyone'
    email: obligatorio si type es 'user' o 'group'.
    """
    try:
        service = _get_service()
        perm: dict[str, Any] = {"role": role, "type": type}
        if type in ("user", "group"):
            if not email:
                return {"ok": False, "action": "drive_share",
                        "error": "Para type 'user'/'group' hace falta 'email'."}
            perm["emailAddress"] = email
        res = service.permissions().create(
            fileId=file_id, body=perm, sendNotificationEmail=notify,
            fields="id, role, type, emailAddress", supportsAllDrives=True,
        ).execute()
        _audit("drive_share", {"file_id": file_id, "role": role, "type": type,
                               "email": email})
        return {"ok": True, "permission": res}
    except Exception as e:
        return _err("drive_share", e)


@mcp.tool()
def drive_list_permissions(file_id: str) -> dict:
    """Lista los permisos (con quién está compartido) de un archivo/carpeta."""
    try:
        service = _get_service()
        res = service.permissions().list(
            fileId=file_id,
            fields="permissions(id, role, type, emailAddress, displayName)",
            supportsAllDrives=True,
        ).execute()
        return {"ok": True, "permissions": res.get("permissions", [])}
    except Exception as e:
        return _err("drive_list_permissions", e)


# --------------------------------------------------------------------------- #
# SUBIR / DESCARGAR
# --------------------------------------------------------------------------- #

@mcp.tool()
def drive_upload_text(name: str, text_content: str, parent_id: str = "root",
                      mime_type: str = "text/plain", convert_to_doc: bool = False) -> dict:
    """Sube un archivo de texto nuevo con el contenido dado.

    Por defecto lo deja EN CRUDO con `mime_type` (p.ej. 'text/markdown' para poder
    trabajarlo luego con IA). Si `convert_to_doc=True`, Drive convierte el contenido
    (idealmente Markdown) en un Documento de Google FORMATEADO, en el servidor y sin
    base64: úsalo solo cuando quieras el documento «bonito»."""
    try:
        service = _get_service()
        media = MediaIoBaseUpload(
            io.BytesIO(text_content.encode("utf-8")), mimetype=mime_type, resumable=False
        )
        body = {"name": name, "parents": [parent_id]}
        if convert_to_doc:
            body["mimeType"] = "application/vnd.google-apps.document"
        meta = service.files().create(
            body=body, media_body=media, fields=FILE_FIELDS, supportsAllDrives=True
        ).execute()
        return {"ok": True, "file": _file(meta)}
    except Exception as e:
        return _err("drive_upload_text", e)


@mcp.tool()
def drive_upload_base64(name: str, base64_content: str, mime_type: str,
                        parent_id: str = "root") -> dict:
    """Sube un archivo binario nuevo a partir de contenido en base64."""
    try:
        service = _get_service()
        data = base64.b64decode(base64_content)
        media = MediaIoBaseUpload(io.BytesIO(data), mimetype=mime_type, resumable=True)
        body = {"name": name, "parents": [parent_id]}
        meta = service.files().create(
            body=body, media_body=media, fields=FILE_FIELDS, supportsAllDrives=True
        ).execute()
        return {"ok": True, "file": _file(meta)}
    except Exception as e:
        return _err("drive_upload_base64", e)


@mcp.tool()
def drive_download_text(file_id: str) -> dict:
    """Descarga el contenido de texto de un archivo. Los documentos de Google
    se exportan a texto plano; el resto se devuelve como UTF-8 si es posible."""
    try:
        service = _get_service()
        meta = service.files().get(
            fileId=file_id, fields="mimeType, name", supportsAllDrives=True
        ).execute()
        mime = meta.get("mimeType", "")
        if mime.startswith("application/vnd.google-apps"):
            export_map = {
                "application/vnd.google-apps.document": "text/plain",
                "application/vnd.google-apps.spreadsheet": "text/csv",
                "application/vnd.google-apps.presentation": "text/plain",
            }
            export_mime = export_map.get(mime, "text/plain")
            data = service.files().export(fileId=file_id, mimeType=export_mime).execute()
            text = data.decode("utf-8", errors="replace") if isinstance(data, bytes) else str(data)
        else:
            request = service.files().get_media(fileId=file_id)
            buf = io.BytesIO()
            downloader = MediaIoBaseDownload(buf, request)
            done = False
            while not done:
                _, done = downloader.next_chunk()
            text = buf.getvalue().decode("utf-8", errors="replace")
        return {"ok": True, "name": meta.get("name"), "content": text}
    except Exception as e:
        return _err("drive_download_text", e)


@mcp.tool()
def drive_read_file(file_id: str, max_chars: int = 200000) -> dict:
    """Lee el CONTENIDO en texto de casi cualquier archivo de Drive: PDF, Word
    (.docx), Excel (.xlsx), PowerPoint (.pptx), texto/Markdown/CSV/JSON, Google
    nativo, e imágenes o PDF escaneado (por OCR). Extrae el texto en el servidor,
    sin base64 ni pasar el binario por el modelo. Recorta a `max_chars`."""
    try:
        svc = _get_service()
        meta = svc.files().get(fileId=file_id, fields="name, mimeType",
                               supportsAllDrives=True).execute()
        name = meta.get("name", "")
        mime = meta.get("mimeType", "")
        ext = name.rsplit(".", 1)[-1].lower() if "." in name else ""

        if mime.startswith("application/vnd.google-apps"):
            return drive_download_text(file_id)

        request = svc.files().get_media(fileId=file_id)
        buf = io.BytesIO()
        dl = MediaIoBaseDownload(buf, request)
        done = False
        while not done:
            _, done = dl.next_chunk()
        data = buf.getvalue()

        text = ""
        if mime == "application/pdf" or ext == "pdf":
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(data))
            text = "\n".join((p.extract_text() or "") for p in reader.pages).strip()
            if not text:
                return drive_ocr(file_id)
        elif ext == "docx" or "wordprocessingml" in mime:
            import docx
            d = docx.Document(io.BytesIO(data))
            text = "\n".join(p.text for p in d.paragraphs)
        elif ext == "xlsx" or "spreadsheetml" in mime:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
            rows = []
            for ws in wb.worksheets:
                rows.append("# " + ws.title)
                for r in ws.iter_rows(values_only=True):
                    rows.append("\t".join("" if c is None else str(c) for c in r))
            text = "\n".join(rows)
        elif ext == "pptx" or "presentationml" in mime:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(data))
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                parts.append("--- Diapositiva %d ---" % i)
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False):
                        parts.append(shape.text_frame.text)
            text = "\n".join(parts)
        elif mime.startswith("text/") or ext in ("txt", "md", "csv", "json", "xml", "html", "log"):
            text = data.decode("utf-8", errors="replace")
        elif mime.startswith("image/"):
            return drive_ocr(file_id)
        else:
            return {"ok": False, "action": "drive_read_file",
                    "error": "Formato no soportado para lectura de texto: " + (mime or ext)}

        return {"ok": True, "name": name, "mimeType": mime,
                "truncated": len(text) > max_chars, "content": text[:max_chars]}
    except Exception as e:
        return _err("drive_read_file", e)


# --------------------------------------------------------------------------- #
# EDICIÓN DE DOCUMENTOS DE GOOGLE (Docs API) — editar el CUERPO de un Doc
# --------------------------------------------------------------------------- #

def _doc_end_index(docs, document_id: str) -> int:
    """Índice de inserción al final del cuerpo (antes del salto de línea final)."""
    doc = docs.documents().get(documentId=document_id).execute()
    content = doc.get("body", {}).get("content", [])
    end = content[-1].get("endIndex", 2) if content else 2
    return end - 1


@mcp.tool()
def drive_doc_get_text(document_id: str) -> dict:
    """Devuelve el texto plano de un Documento de Google (para inspeccionarlo)."""
    try:
        docs = _get_docs()
        doc = docs.documents().get(documentId=document_id).execute()
        out = []
        for el in doc.get("body", {}).get("content", []):
            para = el.get("paragraph")
            if not para:
                continue
            for pe in para.get("elements", []):
                t = pe.get("textRun", {}).get("content")
                if t:
                    out.append(t)
        return {"ok": True, "title": doc.get("title"), "text": "".join(out)}
    except Exception as e:
        return _err("drive_doc_get_text", e)


@mcp.tool()
def drive_doc_append_heading(document_id: str, text: str, level: int = 2) -> dict:
    """Añade un ENCABEZADO (nivel 1-3) al final de un Documento de Google."""
    try:
        docs = _get_docs()
        idx = _doc_end_index(docs, document_id)
        style = {1: "HEADING_1", 2: "HEADING_2", 3: "HEADING_3"}.get(int(level), "HEADING_2")
        reqs = [
            {"insertText": {"location": {"index": idx}, "text": "\n" + text}},
            {"updateParagraphStyle": {
                "range": {"startIndex": idx + 1, "endIndex": idx + 1 + len(text)},
                "paragraphStyle": {"namedStyleType": style},
                "fields": "namedStyleType"}},
        ]
        docs.documents().batchUpdate(documentId=document_id, body={"requests": reqs}).execute()
        return {"ok": True, "document_id": document_id, "heading": text, "level": int(level)}
    except Exception as e:
        return _err("drive_doc_append_heading", e)


@mcp.tool()
def drive_doc_append_text(document_id: str, text: str) -> dict:
    """Añade un PÁRRAFO de texto normal al final de un Documento de Google."""
    try:
        docs = _get_docs()
        idx = _doc_end_index(docs, document_id)
        reqs = [
            {"insertText": {"location": {"index": idx}, "text": "\n" + text}},
            {"updateParagraphStyle": {
                "range": {"startIndex": idx + 1, "endIndex": idx + 1 + len(text)},
                "paragraphStyle": {"namedStyleType": "NORMAL_TEXT"},
                "fields": "namedStyleType"}},
        ]
        docs.documents().batchUpdate(documentId=document_id, body={"requests": reqs}).execute()
        return {"ok": True, "document_id": document_id, "added_chars": len(text)}
    except Exception as e:
        return _err("drive_doc_append_text", e)


@mcp.tool()
def drive_doc_replace_text(document_id: str, find: str, replace_with: str,
                           match_case: bool = True) -> dict:
    """Reemplaza TODAS las apariciones de un texto dentro de un Documento de Google."""
    try:
        docs = _get_docs()
        reqs = [{"replaceAllText": {
            "containsText": {"text": find, "matchCase": bool(match_case)},
            "replaceText": replace_with}}]
        r = docs.documents().batchUpdate(documentId=document_id, body={"requests": reqs}).execute()
        occ = 0
        for rep in r.get("replies", []):
            occ += rep.get("replaceAllText", {}).get("occurrencesChanged", 0) or 0
        return {"ok": True, "document_id": document_id, "replaced": occ}
    except Exception as e:
        return _err("drive_doc_replace_text", e)


# --------------------------------------------------------------------------- #
# PLANTILLAS DE ESCRITOS (copiar plantilla + sustituir marcadores)
# --------------------------------------------------------------------------- #

@mcp.tool()
def drive_doc_from_template(template_id: str, new_name: str,
                            replacements: dict, parent_id: str = "root") -> dict:
    """Crea un documento a partir de una plantilla: copia el Doc `template_id`,
    sustituye los marcadores de `replacements` (p.ej. {"{{cliente}}": "Ana"}) y
    lo guarda como `new_name` en `parent_id`. Ideal para demandas, papeletas, etc."""
    try:
        svc = _get_service()
        meta = svc.files().copy(
            fileId=template_id, body={"name": new_name, "parents": [parent_id]},
            fields=FILE_FIELDS, supportsAllDrives=True,
        ).execute()
        reqs = [{"replaceAllText": {"containsText": {"text": k, "matchCase": True},
                                    "replaceText": str(v)}}
                for k, v in (replacements or {}).items()]
        if reqs:
            _get_docs().documents().batchUpdate(
                documentId=meta["id"], body={"requests": reqs}).execute()
        return {"ok": True, "file": _file(meta), "sustituciones": len(reqs)}
    except Exception as e:
        return _err("drive_doc_from_template", e)


# --------------------------------------------------------------------------- #
# HOJAS DE CÁLCULO (Google Sheets)
# --------------------------------------------------------------------------- #

_sheets_service = None


def _get_sheets():
    global _sheets_service
    if _sheets_service is None:
        _sheets_service = build("sheets", "v4", credentials=_build_creds(), cache_discovery=False)
    return _sheets_service


@mcp.tool()
def drive_sheet_read(spreadsheet_id: str, range_a1: str) -> dict:
    """Lee un rango en notación A1 (p.ej. 'Hoja1!A1:D20') de una hoja de cálculo."""
    try:
        r = _get_sheets().spreadsheets().values().get(
            spreadsheetId=spreadsheet_id, range=range_a1).execute()
        return {"ok": True, "range": r.get("range"), "values": r.get("values", [])}
    except Exception as e:
        return _err("drive_sheet_read", e)


@mcp.tool()
def drive_sheet_write(spreadsheet_id: str, range_a1: str, values: list) -> dict:
    """Escribe una matriz de valores (lista de listas) en un rango A1."""
    try:
        r = _get_sheets().spreadsheets().values().update(
            spreadsheetId=spreadsheet_id, range=range_a1,
            valueInputOption="USER_ENTERED", body={"values": values}).execute()
        return {"ok": True, "celdas_actualizadas": r.get("updatedCells")}
    except Exception as e:
        return _err("drive_sheet_write", e)


@mcp.tool()
def drive_sheet_append_row(spreadsheet_id: str, range_a1: str, values: list) -> dict:
    """Añade una fila (o varias) al final de la tabla que contiene `range_a1`."""
    try:
        rows = values if (values and isinstance(values[0], list)) else [values]
        r = _get_sheets().spreadsheets().values().append(
            spreadsheetId=spreadsheet_id, range=range_a1,
            valueInputOption="USER_ENTERED", insertDataOption="INSERT_ROWS",
            body={"values": rows}).execute()
        return {"ok": True, "actualizacion": r.get("updates", {})}
    except Exception as e:
        return _err("drive_sheet_append_row", e)


@mcp.tool()
def drive_create_sheet(name: str, parent_id: str = "root",
                       csv_content: str = "", delimiter: str = ",") -> dict:
    """Crea una NUEVA Hoja de cálculo de Google (nativa) en Drive.

    - Sin `csv_content`: crea una hoja VACÍA.
    - Con `csv_content`: sube ese CSV/TSV y Drive lo CONVIERTE a Hoja nativa
      en el servidor. El contenido viaja como media 'text/csv' (o
      'text/tab-separated-values' si `delimiter` es tabulador) y el destino
      'application/vnd.google-apps.spreadsheet' va en los METADATOS del
      fichero. No se pone el mime de Google en el media: la API lo rechaza
      con "Invalid MIME type provided for the uploaded content".

    Devuelve el fichero creado; su `id` es el `spreadsheet_id` que luego
    usan drive_sheet_read / drive_sheet_write / drive_sheet_append_row.
    """
    try:
        service = _get_service()
        body = {"name": name,
                "mimeType": "application/vnd.google-apps.spreadsheet",
                "parents": [parent_id]}
        if csv_content:
            media_mime = ("text/tab-separated-values"
                          if delimiter == "\t" else "text/csv")
            media = MediaIoBaseUpload(
                io.BytesIO(csv_content.encode("utf-8")),
                mimetype=media_mime, resumable=False)
            meta = service.files().create(
                body=body, media_body=media, fields=FILE_FIELDS,
                supportsAllDrives=True).execute()
        else:
            meta = service.files().create(
                body=body, fields=FILE_FIELDS,
                supportsAllDrives=True).execute()
        return {"ok": True, "file": _file(meta)}
    except Exception as e:
        return _err("drive_create_sheet", e)


# --------------------------------------------------------------------------- #
# NOVEDADES Y AUDITORÍA
# --------------------------------------------------------------------------- #

@mcp.tool()
def drive_novedades(since: str, page_size: int = 50) -> dict:
    """Lista lo creado o modificado desde una fecha (RFC3339, p.ej.
    '2026-07-25T00:00:00Z'). Sirve para saber 'qué hay nuevo' en el Drive."""
    try:
        q = (f"(createdTime > '{since}' or modifiedTime > '{since}') "
             f"and trashed = false")
        res = _get_service().files().list(
            q=q, pageSize=max(1, min(page_size, 1000)), orderBy="modifiedTime desc",
            fields=f"files({FILE_FIELDS})", supportsAllDrives=True,
            includeItemsFromAllDrives=True).execute()
        files = [_file(f) for f in res.get("files", [])]
        return {"ok": True, "desde": since, "count": len(files), "files": files}
    except Exception as e:
        return _err("drive_novedades", e)


@mcp.tool()
def drive_compartidos_externos(page_size: int = 100) -> dict:
    """Lista archivos compartidos por enlace público (cualquiera con el enlace).
    Auditoría de seguridad: útil por los datos de clientes en 01_JURIDICO."""
    try:
        q = ("(visibility = 'anyoneWithLink' or visibility = 'anyoneCanFind') "
             "and trashed = false")
        res = _get_service().files().list(
            q=q, pageSize=max(1, min(page_size, 1000)),
            fields=f"files({FILE_FIELDS})", supportsAllDrives=True,
            includeItemsFromAllDrives=True).execute()
        files = [_file(f) for f in res.get("files", [])]
        return {"ok": True, "count": len(files), "files": files}
    except Exception as e:
        return _err("drive_compartidos_externos", e)


@mcp.tool()
def drive_search_content(text: str, page_size: int = 25) -> dict:
    """Busca por el CONTENIDO (texto interno) de los archivos, no solo el título."""
    try:
        safe = text.replace("'", "\\'")
        q = f"fullText contains '{safe}' and trashed = false"
        res = _get_service().files().list(
            q=q, pageSize=max(1, min(page_size, 1000)),
            fields=f"files({FILE_FIELDS})", supportsAllDrives=True,
            includeItemsFromAllDrives=True).execute()
        files = [_file(f) for f in res.get("files", [])]
        return {"ok": True, "count": len(files), "files": files}
    except Exception as e:
        return _err("drive_search_content", e)


# --------------------------------------------------------------------------- #
# DOCS: viñetas, salto de página, y exportación a PDF
# --------------------------------------------------------------------------- #

@mcp.tool()
def drive_doc_append_bullets(document_id: str, items: list) -> dict:
    """Añade una lista con viñetas (una por cada elemento de `items`)."""
    try:
        docs = _get_docs()
        idx = _doc_end_index(docs, document_id)
        joined = "\n".join(items)
        reqs = [
            {"insertText": {"location": {"index": idx}, "text": "\n" + joined}},
            {"createParagraphBullets": {
                "range": {"startIndex": idx + 1, "endIndex": idx + 1 + len(joined)},
                "bulletPreset": "BULLET_DISC_CIRCLE_SQUARE"}},
        ]
        docs.documents().batchUpdate(documentId=document_id, body={"requests": reqs}).execute()
        return {"ok": True, "document_id": document_id, "items": len(items)}
    except Exception as e:
        return _err("drive_doc_append_bullets", e)


@mcp.tool()
def drive_doc_insert_page_break(document_id: str) -> dict:
    """Añade un salto de página al final del documento."""
    try:
        docs = _get_docs()
        idx = _doc_end_index(docs, document_id)
        reqs = [{"insertPageBreak": {"location": {"index": idx}}}]
        docs.documents().batchUpdate(documentId=document_id, body={"requests": reqs}).execute()
        return {"ok": True, "document_id": document_id}
    except Exception as e:
        return _err("drive_doc_insert_page_break", e)


@mcp.tool()
def drive_export_pdf(file_id: str, name: str, parent_id: str = "root") -> dict:
    """Exporta un documento de Google (Doc/Sheet/Slide) a PDF y lo guarda en Drive."""
    try:
        svc = _get_service()
        data = svc.files().export(fileId=file_id, mimeType="application/pdf").execute()
        media = MediaIoBaseUpload(io.BytesIO(data), mimetype="application/pdf", resumable=True)
        meta = svc.files().create(
            body={"name": name, "parents": [parent_id]}, media_body=media,
            fields=FILE_FIELDS, supportsAllDrives=True).execute()
        return {"ok": True, "file": _file(meta)}
    except Exception as e:
        return _err("drive_export_pdf", e)


# --------------------------------------------------------------------------- #
# GESTIÓN AVANZADA (papelera, lotes, accesos directos, color, estrella)
# --------------------------------------------------------------------------- #

@mcp.tool()
def drive_list_trash(page_size: int = 100) -> dict:
    """Lista el contenido de la papelera."""
    try:
        res = _get_service().files().list(
            q="trashed = true", pageSize=max(1, min(page_size, 1000)),
            fields=f"files({FILE_FIELDS})", supportsAllDrives=True,
            includeItemsFromAllDrives=True).execute()
        files = [_file(f) for f in res.get("files", [])]
        return {"ok": True, "count": len(files), "files": files}
    except Exception as e:
        return _err("drive_list_trash", e)


@mcp.tool()
def drive_empty_trash(confirm: bool = False) -> dict:
    """Vacía la papelera DEFINITIVAMENTE (no reversible). Requiere confirm=True."""
    if not confirm:
        return {"ok": False, "action": "drive_empty_trash",
                "error": "Bloqueado. Vaciar la papelera es irreversible; vuelve a "
                         "llamar con confirm=true solo si el usuario lo pide expresamente."}
    try:
        _get_service().files().emptyTrash().execute()
        _audit("drive_empty_trash", {"papelera": "vaciada"})
        return {"ok": True, "papelera": "vaciada"}
    except Exception as e:
        return _err("drive_empty_trash", e)


@mcp.tool()
def drive_batch_move(file_ids: list, new_parent_id: str, dry_run: bool = False) -> dict:
    """Mueve varios archivos/carpetas a `new_parent_id` de una vez.
    Con dry_run=True informa de lo que haria sin mover nada."""
    results = []
    try:
        svc = _get_service()
        for fid in file_ids:
            try:
                cur = svc.files().get(fileId=fid, fields="parents, name", supportsAllDrives=True).execute()
                prev = ",".join(cur.get("parents", []))
                if dry_run:
                    results.append({"id": fid, "name": cur.get("name"),
                                    "dry_run": True, "hacia": new_parent_id})
                    continue
                svc.files().update(fileId=fid, addParents=new_parent_id,
                                   removeParents=prev, fields="id, name",
                                   supportsAllDrives=True).execute()
                _audit("drive_batch_move", {"file_id": fid, "name": cur.get("name"),
                                            "hacia": new_parent_id})
                results.append({"id": fid, "ok": True})
            except Exception as ie:
                results.append({"id": fid, "ok": False, "error": str(ie)})
        return {"ok": True, "dry_run": dry_run, "resultados": results}
    except Exception as e:
        return _err("drive_batch_move", e)


@mcp.tool()
def drive_create_shortcut(target_id: str, name: str, parent_id: str = "root") -> dict:
    """Crea un acceso directo a `target_id` dentro de `parent_id`."""
    try:
        body = {"name": name, "mimeType": "application/vnd.google-apps.shortcut",
                "parents": [parent_id], "shortcutDetails": {"targetId": target_id}}
        meta = _get_service().files().create(
            body=body, fields=FILE_FIELDS, supportsAllDrives=True).execute()
        return {"ok": True, "file": _file(meta)}
    except Exception as e:
        return _err("drive_create_shortcut", e)


@mcp.tool()
def drive_set_star(file_id: str, starred: bool = True) -> dict:
    """Marca o desmarca un archivo como Destacado."""
    try:
        meta = _get_service().files().update(
            fileId=file_id, body={"starred": bool(starred)},
            fields=FILE_FIELDS, supportsAllDrives=True).execute()
        return {"ok": True, "file": _file(meta)}
    except Exception as e:
        return _err("drive_set_star", e)


@mcp.tool()
def drive_set_folder_color(folder_id: str, color_rgb: str) -> dict:
    """Cambia el color de una carpeta (hex, p.ej. '#d50000' rojo, '#3f51b5' azul)."""
    try:
        meta = _get_service().files().update(
            fileId=folder_id, body={"folderColorRgb": color_rgb},
            fields=FILE_FIELDS, supportsAllDrives=True).execute()
        return {"ok": True, "file": _file(meta)}
    except Exception as e:
        return _err("drive_set_folder_color", e)


# --------------------------------------------------------------------------- #
# COMENTARIOS, REVISIONES, OCR Y CUOTA
# --------------------------------------------------------------------------- #

@mcp.tool()
def drive_list_comments(file_id: str) -> dict:
    """Lista los comentarios de un archivo (Doc, Sheet...)."""
    try:
        res = _get_service().comments().list(
            fileId=file_id,
            fields="comments(id, content, resolved, createdTime, author/displayName)").execute()
        return {"ok": True, "comments": res.get("comments", [])}
    except Exception as e:
        return _err("drive_list_comments", e)


@mcp.tool()
def drive_add_comment(file_id: str, content: str) -> dict:
    """Añade un comentario a un archivo."""
    try:
        res = _get_service().comments().create(
            fileId=file_id, body={"content": content},
            fields="id, content, createdTime").execute()
        return {"ok": True, "comment": res}
    except Exception as e:
        return _err("drive_add_comment", e)


@mcp.tool()
def drive_list_revisions(file_id: str) -> dict:
    """Lista el historial de versiones de un archivo."""
    try:
        res = _get_service().revisions().list(
            fileId=file_id,
            fields="revisions(id, modifiedTime, size, lastModifyingUser/displayName)").execute()
        return {"ok": True, "revisions": res.get("revisions", [])}
    except Exception as e:
        return _err("drive_list_revisions", e)


@mcp.tool()
def drive_ocr(file_id: str, language: str = "es", new_name: str = "OCR") -> dict:
    """Extrae el texto de un PDF o imagen mediante OCR (crea un Doc con el texto
    reconocido y devuelve ese texto). `language` en ISO (es, en...)."""
    try:
        svc = _get_service()
        doc = svc.files().copy(
            fileId=file_id, ocrLanguage=language,
            body={"name": new_name, "mimeType": "application/vnd.google-apps.document"},
            fields="id, name", supportsAllDrives=True).execute()
        data = svc.files().export(fileId=doc["id"], mimeType="text/plain").execute()
        text = data.decode("utf-8", errors="replace") if isinstance(data, bytes) else str(data)
        return {"ok": True, "doc_id": doc["id"], "text": text}
    except Exception as e:
        return _err("drive_ocr", e)


@mcp.tool()
def drive_storage_quota() -> dict:
    """Devuelve el uso de almacenamiento de la cuenta de Drive."""
    try:
        r = _get_service().about().get(fields="storageQuota, user").execute()
        return {"ok": True, "storageQuota": r.get("storageQuota"), "user": r.get("user")}
    except Exception as e:
        return _err("drive_storage_quota", e)


# --------------------------------------------------------------------------- #
# GOBERNANZA: auditoria, ciclo de permisos y remediacion de compartidos
# --------------------------------------------------------------------------- #

@mcp.tool()
def drive_audit_log(n: int = 50) -> dict:
    """Devuelve las ultimas n entradas del registro de auditoria (acciones que
    modifican Drive: mover, renombrar, papelera, borrar, compartir, permisos)."""
    try:
        if not os.path.exists(AUDIT_PATH):
            return {"ok": True, "count": 0, "entries": [], "nota": "sin registro todavia"}
        with open(AUDIT_PATH, encoding="utf-8") as f:
            lines = [l for l in f.readlines() if l.strip()]
        sel = lines[-max(1, min(n, 1000)):]
        return {"ok": True, "count": len(sel), "entries": [_json.loads(l) for l in sel]}
    except Exception as e:
        return _err("drive_audit_log", e)


@mcp.tool()
def drive_remove_permission(file_id: str, permission_id: str) -> dict:
    """Revoca (elimina) un permiso concreto de un archivo/carpeta.
    Usa drive_list_permissions para ver los ids de permiso."""
    try:
        _get_service().permissions().delete(
            fileId=file_id, permissionId=permission_id, supportsAllDrives=True
        ).execute()
        _audit("drive_remove_permission", {"file_id": file_id, "permission_id": permission_id})
        return {"ok": True, "file_id": file_id, "permiso_eliminado": permission_id}
    except Exception as e:
        return _err("drive_remove_permission", e)


@mcp.tool()
def drive_update_permission(file_id: str, permission_id: str, role: str) -> dict:
    """Cambia el rol de un permiso existente (reader | commenter | writer)."""
    try:
        res = _get_service().permissions().update(
            fileId=file_id, permissionId=permission_id, body={"role": role},
            fields="id, role, type, emailAddress", supportsAllDrives=True,
        ).execute()
        _audit("drive_update_permission", {"file_id": file_id,
               "permission_id": permission_id, "role": role})
        return {"ok": True, "permission": res}
    except Exception as e:
        return _err("drive_update_permission", e)


@mcp.tool()
def drive_remediar_externos(page_size: int = 100, dry_run: bool = True) -> dict:
    """Encuentra ficheros compartidos por enlace publico (cualquiera con el enlace)
    y revoca ese permiso. Con dry_run=True (por defecto) solo informa de lo que
    quitaria; con dry_run=False elimina los permisos de tipo 'anyone'. Auditoria
    de seguridad para datos de clientes."""
    try:
        svc = _get_service()
        q = ("(visibility = 'anyoneWithLink' or visibility = 'anyoneCanFind') "
             "and trashed = false")
        res = svc.files().list(q=q, pageSize=max(1, min(page_size, 1000)),
                               fields=f"files({FILE_FIELDS})", supportsAllDrives=True,
                               includeItemsFromAllDrives=True).execute()
        # Recoger TODOS los ficheros publicos (paginando)
        files=[]; token=None
        while True:
            r = svc.files().list(q=q, pageSize=1000, pageToken=token,
                fields=f"nextPageToken, files({FILE_FIELDS})",
                supportsAllDrives=True, includeItemsFromAllDrives=True).execute()
            files.extend(r.get("files", []))
            token = r.get("nextPageToken")
            if not token:
                break
        acciones = []
        omitidos = 0
        for f in files:
            try:
                perms = svc.permissions().list(
                    fileId=f["id"], fields="permissions(id, type, role)",
                    supportsAllDrives=True).execute().get("permissions", [])
                for pm in perms:
                    if pm.get("type") == "anyone":
                        if dry_run:
                            acciones.append({"file_id": f["id"], "name": f.get("name"),
                                             "mimeType": f.get("mimeType"),
                                             "quitaria_permiso": pm.get("id"),
                                             "role": pm.get("role")})
                        else:
                            svc.permissions().delete(fileId=f["id"], permissionId=pm["id"],
                                                     supportsAllDrives=True).execute()
                            _audit("drive_remediar_externos",
                                   {"file_id": f["id"], "name": f.get("name"),
                                    "permiso_eliminado": pm.get("id")})
                            acciones.append({"file_id": f["id"], "name": f.get("name"),
                                             "permiso_eliminado": pm.get("id"),
                                             "role": pm.get("role")})
            except Exception:
                omitidos += 1
                continue
        return {"ok": True, "dry_run": dry_run, "num": len(acciones),
                "omitidos_sin_permiso": omitidos, "acciones": acciones}
    except Exception as e:
        return _err("drive_remediar_externos", e)



# --------------------------------------------------------------------------- #
# GMAIL (lectura de correo y volcado de adjuntos a Drive)
# --------------------------------------------------------------------------- #

_gmail_service = None

def _get_gmail():
    global _gmail_service
    if _gmail_service is None:
        _gmail_service = build("gmail", "v1", credentials=_build_creds(), cache_discovery=False)
    return _gmail_service


def _gmail_walk(payload):
    stack = [payload]
    while stack:
        p = stack.pop()
        for c in (p.get("parts") or []):
            stack.append(c)
        yield p


@mcp.tool()
def gmail_buscar(query: str, max_results: int = 20) -> dict:
    """Busca correos en Gmail con la sintaxis de Gmail. Ejemplos de `query`:
    'has:attachment newer_than:7d', 'from:cliente@dominio.com', 'subject:nomina'."""
    try:
        svc = _get_gmail()
        res = svc.users().messages().list(userId="me", q=query,
              maxResults=max(1, min(max_results, 100))).execute()
        out = []
        for m in res.get("messages", []):
            full = svc.users().messages().get(userId="me", id=m["id"], format="metadata",
                   metadataHeaders=["From", "Subject", "Date"]).execute()
            h = {x["name"]: x["value"] for x in full.get("payload", {}).get("headers", [])}
            out.append({"id": m["id"], "from": h.get("From"), "subject": h.get("Subject"),
                        "date": h.get("Date"), "snippet": full.get("snippet", "")[:160]})
        return {"ok": True, "count": len(out), "messages": out}
    except Exception as e:
        return _err("gmail_buscar", e)


@mcp.tool()
def gmail_leer(message_id: str, max_chars: int = 20000) -> dict:
    """Lee un correo de Gmail: remitente, asunto, fecha, texto y lista de adjuntos."""
    try:
        import base64 as _b64
        svc = _get_gmail()
        m = svc.users().messages().get(userId="me", id=message_id, format="full").execute()
        payload = m.get("payload", {})
        h = {x["name"]: x["value"] for x in payload.get("headers", [])}
        text = ""
        adjs = []
        for p in _gmail_walk(payload):
            fn = p.get("filename")
            body = p.get("body", {})
            if fn:
                adjs.append({"filename": fn, "attachmentId": body.get("attachmentId"),
                             "size": body.get("size")})
            elif p.get("mimeType") == "text/plain" and body.get("data"):
                text += _b64.urlsafe_b64decode(body["data"]).decode("utf-8", "replace")
        return {"ok": True, "from": h.get("From"), "subject": h.get("Subject"),
                "date": h.get("Date"), "texto": text[:max_chars], "adjuntos": adjs}
    except Exception as e:
        return _err("gmail_leer", e)


@mcp.tool()
def gmail_guardar_adjuntos(message_id: str, parent_id: str = "root") -> dict:
    """Descarga los adjuntos de un correo de Gmail y los sube a una carpeta de Drive."""
    try:
        import base64 as _b64
        svc = _get_gmail()
        drv = _get_service()
        m = svc.users().messages().get(userId="me", id=message_id, format="full").execute()
        saved = []
        for p in _gmail_walk(m.get("payload", {})):
            fn = p.get("filename")
            body = p.get("body", {})
            if fn and body.get("attachmentId"):
                att = svc.users().messages().attachments().get(
                    userId="me", messageId=message_id, id=body["attachmentId"]).execute()
                data = _b64.urlsafe_b64decode(att["data"])
                media = MediaIoBaseUpload(io.BytesIO(data),
                        mimetype="application/octet-stream", resumable=False)
                r = drv.files().create(body={"name": fn, "parents": [parent_id]},
                    media_body=media, fields="id, name").execute()
                _audit("gmail_guardar_adjuntos", {"message_id": message_id, "file": fn,
                       "drive_id": r["id"]})
                saved.append({"filename": fn, "drive_id": r["id"]})
        return {"ok": True, "guardados": len(saved), "adjuntos": saved}
    except Exception as e:
        return _err("gmail_guardar_adjuntos", e)



# --------------------------------------------------------------------------- #
# CALENDAR / MEET (agenda y reuniones con videollamada de Google Meet)
# --------------------------------------------------------------------------- #

_cal_service = None

def _get_cal():
    global _cal_service
    if _cal_service is None:
        _cal_service = build("calendar", "v3", credentials=_build_creds(), cache_discovery=False)
    return _cal_service


@mcp.tool()
def calendar_listar_eventos(time_min: Optional[str] = None, time_max: Optional[str] = None,
                            max_results: int = 20, calendar_id: str = "primary") -> dict:
    """Lista eventos del calendario. Fechas en RFC3339 (2026-08-05T00:00:00Z). Si no se
    da time_min, usa ahora (próximos eventos)."""
    try:
        import datetime as _dt
        svc = _get_cal()
        if not time_min:
            time_min = _dt.datetime.utcnow().isoformat() + "Z"
        params = dict(calendarId=calendar_id, timeMin=time_min, singleEvents=True,
                      orderBy="startTime", maxResults=max(1, min(max_results, 100)))
        if time_max:
            params["timeMax"] = time_max
        res = svc.events().list(**params).execute()
        out = []
        for e in res.get("items", []):
            st = e.get("start", {}); en = e.get("end", {})
            out.append({"id": e.get("id"), "summary": e.get("summary"),
                        "inicio": st.get("dateTime") or st.get("date"),
                        "fin": en.get("dateTime") or en.get("date"),
                        "meet": e.get("hangoutLink"),
                        "asistentes": [a.get("email") for a in e.get("attendees", [])],
                        "link": e.get("htmlLink")})
        return {"ok": True, "count": len(out), "eventos": out}
    except Exception as e:
        return _err("calendar_listar_eventos", e)


@mcp.tool()
def calendar_crear_evento(summary: str, start: str, end: str, description: str = "",
                          attendees: Optional[list] = None, con_meet: bool = True,
                          calendar_id: str = "primary") -> dict:
    """Crea un evento en el calendario. start/end en RFC3339 con zona
    (2026-08-05T10:00:00+01:00). con_meet=True añade enlace de Google Meet.
    attendees = lista de emails (se les envía invitación)."""
    try:
        import uuid as _uuid
        svc = _get_cal()
        body = {"summary": summary, "description": description,
                "start": {"dateTime": start}, "end": {"dateTime": end}}
        if attendees:
            body["attendees"] = [{"email": a} for a in attendees]
        params = dict(calendarId=calendar_id, body=body, sendUpdates="all")
        if con_meet:
            body["conferenceData"] = {"createRequest": {"requestId": str(_uuid.uuid4()),
                "conferenceSolutionKey": {"type": "hangoutsMeet"}}}
            params["conferenceDataVersion"] = 1
        e = svc.events().insert(**params).execute()
        _audit("calendar_crear_evento", {"summary": summary, "start": start})
        return {"ok": True, "id": e.get("id"), "meet": e.get("hangoutLink"), "link": e.get("htmlLink")}
    except Exception as e:
        return _err("calendar_crear_evento", e)


@mcp.tool()
def calendar_buscar(query: str, max_results: int = 20, calendar_id: str = "primary") -> dict:
    """Busca eventos por texto (nombre, asistente, descripción)."""
    try:
        svc = _get_cal()
        res = svc.events().list(calendarId=calendar_id, q=query, singleEvents=True,
              orderBy="startTime", maxResults=max(1, min(max_results, 100))).execute()
        out = [{"id": e.get("id"), "summary": e.get("summary"),
                "inicio": (e.get("start") or {}).get("dateTime") or (e.get("start") or {}).get("date"),
                "link": e.get("htmlLink")} for e in res.get("items", [])]
        return {"ok": True, "count": len(out), "eventos": out}
    except Exception as e:
        return _err("calendar_buscar", e)


@mcp.tool()
def calendar_disponibilidad(time_min: str, time_max: str, calendar_id: str = "primary") -> dict:
    """Devuelve los tramos OCUPADOS entre dos fechas (RFC3339), para localizar huecos."""
    try:
        svc = _get_cal()
        res = svc.freebusy().query(body={"timeMin": time_min, "timeMax": time_max,
              "items": [{"id": calendar_id}]}).execute()
        busy = res.get("calendars", {}).get(calendar_id, {}).get("busy", [])
        return {"ok": True, "ocupado": busy}
    except Exception as e:
        return _err("calendar_disponibilidad", e)


# ==========================================================================
# BLOQUE: Gmail de ESCRITURA (enviar, responder, borradores, etiquetas)
# ==========================================================================


@mcp.tool()
def gmail_enviar(para: str, asunto: str, cuerpo: str, cc: Optional[str] = None, cco: Optional[str] = None, adjuntos_drive: Optional[list] = None):
    """Envia un correo desde Gmail. Compone un mensaje MIME con destinatario, asunto y cuerpo de texto plano; admite copia (cc). Si se indican fileIds de Google Drive en adjuntos_drive, descarga cada archivo y lo adjunta al correo. Devuelve el id del mensaje enviado."""
    try:
        from email.mime.text import MIMEText
        from email.mime.multipart import MIMEMultipart
        from email.mime.base import MIMEBase
        from email import encoders
        svc = _get_gmail()
        msg = MIMEMultipart()
        msg['To'] = para
        msg['Subject'] = asunto
        if cc:
            msg['Cc'] = cc
        if cco:
            msg['Bcc'] = cco
        msg.attach(MIMEText(cuerpo, 'plain', 'utf-8'))
        if adjuntos_drive:
            drive = _get_service()
            for fid in adjuntos_drive:
                meta = drive.files().get(fileId=fid, fields='id, name, mimeType').execute()
                nombre = meta.get('name', str(fid))
                buf = io.BytesIO()
                downloader = MediaIoBaseDownload(buf, drive.files().get_media(fileId=fid))
                done = False
                while not done:
                    _progreso, done = downloader.next_chunk()
                buf.seek(0)
                part = MIMEBase('application', 'octet-stream')
                part.set_payload(buf.read())
                encoders.encode_base64(part)
                part.add_header('Content-Disposition', 'attachment', filename=nombre)
                msg.attach(part)
        raw = base64.urlsafe_b64encode(msg.as_bytes()).decode()
        enviado = svc.users().messages().send(userId='me', body={'raw': raw}).execute()
        _audit("gmail_enviar", {"para": para, "asunto": asunto, "cc": cc, "adjuntos": adjuntos_drive, "id": enviado.get('id')})
        return {"ok": True, "id": enviado.get('id'), "threadId": enviado.get('threadId')}
    except Exception as e:
        return _err("gmail_enviar", e)


@mcp.tool()
def gmail_responder(message_id: str, texto: str):
    """Responde a un correo existente dentro de su MISMO hilo. Obtiene el mensaje original para extraer el remitente, el asunto y las cabeceras de referencia; compone la respuesta con destinatario el remitente original, asunto con prefijo 'Re:' y cabeceras In-Reply-To y References; y la envia en el mismo threadId. Devuelve el id del mensaje enviado."""
    try:
        from email.mime.text import MIMEText
        svc = _get_gmail()
        original = svc.users().messages().get(
            userId='me', id=message_id, format='metadata',
            metadataHeaders=['From', 'Subject', 'Message-Id', 'References']).execute()
        thread_id = original.get('threadId')
        headers = original.get('payload', {}).get('headers', [])

        def _h(name):
            for h in headers:
                if h.get('name', '').lower() == name.lower():
                    return h.get('value')
            return None

        remitente = _h('From') or ''
        asunto_orig = _h('Subject') or ''
        msg_id_hdr = _h('Message-Id')
        refs = _h('References')
        asunto = asunto_orig if asunto_orig.lower().startswith('re:') else 'Re: ' + asunto_orig
        msg = MIMEText(texto, 'plain', 'utf-8')
        msg['To'] = remitente
        msg['Subject'] = asunto
        if msg_id_hdr:
            msg['In-Reply-To'] = msg_id_hdr
            msg['References'] = (refs + ' ' + msg_id_hdr) if refs else msg_id_hdr
        raw = base64.urlsafe_b64encode(msg.as_bytes()).decode()
        enviado = svc.users().messages().send(userId='me', body={'raw': raw, 'threadId': thread_id}).execute()
        _audit("gmail_responder", {"message_id": message_id, "para": remitente, "asunto": asunto, "id": enviado.get('id')})
        return {"ok": True, "id": enviado.get('id'), "threadId": enviado.get('threadId')}
    except Exception as e:
        return _err("gmail_responder", e)


@mcp.tool()
def gmail_crear_borrador(para: str, asunto: str, cuerpo: str):
    """Crea un borrador de correo en Gmail SIN enviarlo. Compone un mensaje MIME con destinatario, asunto y cuerpo de texto plano y lo guarda como borrador. Devuelve el id del borrador creado."""
    try:
        from email.mime.text import MIMEText
        svc = _get_gmail()
        msg = MIMEText(cuerpo, 'plain', 'utf-8')
        msg['To'] = para
        msg['Subject'] = asunto
        raw = base64.urlsafe_b64encode(msg.as_bytes()).decode()
        borrador = svc.users().drafts().create(userId='me', body={'message': {'raw': raw}}).execute()
        _audit("gmail_crear_borrador", {"para": para, "asunto": asunto, "id": borrador.get('id')})
        return {"ok": True, "id": borrador.get('id')}
    except Exception as e:
        return _err("gmail_crear_borrador", e)


@mcp.tool()
def gmail_etiquetar(message_id: str, add_labels: Optional[list] = None, remove_labels: Optional[list] = None):
    """Modifica las etiquetas de un mensaje de Gmail. Anade las etiquetas indicadas en add_labels y elimina las de remove_labels (identificadas por sus labelIds). Devuelve el mensaje con sus etiquetas actualizadas."""
    try:
        svc = _get_gmail()
        body = {}
        if add_labels:
            body['addLabelIds'] = add_labels
        if remove_labels:
            body['removeLabelIds'] = remove_labels
        res = svc.users().messages().modify(userId='me', id=message_id, body=body).execute()
        _audit("gmail_etiquetar", {"message_id": message_id, "add": add_labels, "remove": remove_labels})
        return {"ok": True, "id": res.get('id'), "labelIds": res.get('labelIds')}
    except Exception as e:
        return _err("gmail_etiquetar", e)


# =====================================================================
# BLOQUE: Calendar (Google Calendar v3)
# Se APPENDEA a server.py. No añade imports ni redefine nada.
# Usa helpers ya existentes: _get_cal(), _err(), _audit().
# =====================================================================


@mcp.tool()
def calendar_actualizar_evento(
    event_id: str,
    summary: Optional[str] = None,
    start: Optional[str] = None,
    end: Optional[str] = None,
    description: Optional[str] = None,
    calendar_id: str = "primary",
):
    """Actualiza (patch) un evento del calendario con los campos no nulos.

    start/end se envían como {'dateTime': valor} en formato RFC3339.
    Notifica a todos los asistentes (sendUpdates='all'). Devuelve id,
    summary y enlace del evento actualizado.
    """
    try:
        cal = _get_cal()
        body = {}
        if summary is not None:
            body["summary"] = summary
        if description is not None:
            body["description"] = description
        if start is not None:
            body["start"] = {"dateTime": start}
        if end is not None:
            body["end"] = {"dateTime": end}
        ev = (
            cal.events()
            .patch(
                calendarId=calendar_id,
                eventId=event_id,
                body=body,
                sendUpdates="all",
            )
            .execute()
        )
        _audit(
            "calendar_actualizar_evento",
            f"calendar={calendar_id} event={event_id} campos={list(body.keys())}",
        )
        return {
            "id": ev.get("id"),
            "summary": ev.get("summary"),
            "link": ev.get("htmlLink"),
        }
    except Exception as e:
        return _err("calendar_actualizar_evento", e)


@mcp.tool()
def calendar_cancelar_evento(event_id: str, calendar_id: str = "primary"):
    """Cancela (borra) un evento avisando a los asistentes.

    Ejecuta events().delete con sendUpdates='all' para notificar la
    cancelación. Devuelve un indicador de resultado.
    """
    try:
        cal = _get_cal()
        cal.events().delete(
            calendarId=calendar_id,
            eventId=event_id,
            sendUpdates="all",
        ).execute()
        _audit(
            "calendar_cancelar_evento",
            f"calendar={calendar_id} event={event_id}",
        )
        return {"ok": True, "event_id": event_id}
    except Exception as e:
        return _err("calendar_cancelar_evento", e)


@mcp.tool()
def calendar_responder_invitacion(
    event_id: str,
    respuesta: str,
    calendar_id: str = "primary",
):
    """Responde a una invitación fijando el responseStatus del asistente.

    'respuesta' debe ser 'accepted', 'declined' o 'tentative'. Localiza en
    los asistentes el que tenga self=True (o el email del propio calendario
    vía getProfile) y hace patch de su responseStatus. Devuelve ok.
    """
    try:
        if respuesta not in ("accepted", "declined", "tentative"):
            raise ValueError(
                "respuesta debe ser 'accepted', 'declined' o 'tentative'"
            )
        cal = _get_cal()
        ev = cal.events().get(
            calendarId=calendar_id, eventId=event_id
        ).execute()
        attendees = ev.get("attendees", []) or []
        # Determinar el email propio.
        self_email = None
        for a in attendees:
            if a.get("self"):
                self_email = a.get("email")
                break
        if self_email is None:
            try:
                prof = cal.calendarList().get(
                    calendarId=calendar_id
                ).execute()
                self_email = prof.get("id")
            except Exception:
                self_email = calendar_id if calendar_id != "primary" else None
        encontrado = False
        for a in attendees:
            if a.get("self") or (
                self_email is not None and a.get("email") == self_email
            ):
                a["responseStatus"] = respuesta
                encontrado = True
                break
        if not encontrado:
            if self_email is None:
                raise ValueError(
                    "No se pudo identificar al asistente propio en el evento"
                )
            attendees.append(
                {"email": self_email, "responseStatus": respuesta}
            )
        ev2 = (
            cal.events()
            .patch(
                calendarId=calendar_id,
                eventId=event_id,
                body={"attendees": attendees},
                sendUpdates="all",
            )
            .execute()
        )
        _audit(
            "calendar_responder_invitacion",
            f"calendar={calendar_id} event={event_id} respuesta={respuesta}",
        )
        return {"ok": True, "event_id": ev2.get("id"), "respuesta": respuesta}
    except Exception as e:
        return _err("calendar_responder_invitacion", e)


@mcp.tool()
def calendar_anadir_asistentes(
    event_id: str,
    emails: list,
    calendar_id: str = "primary",
):
    """Añade asistentes nuevos a un evento y notifica el cambio.

    Obtiene el evento, agrega a la lista de asistentes los emails que aún no
    figuren, y hace patch con sendUpdates='all'. Devuelve la lista final de
    asistentes.
    """
    try:
        cal = _get_cal()
        ev = cal.events().get(
            calendarId=calendar_id, eventId=event_id
        ).execute()
        attendees = ev.get("attendees", []) or []
        existentes = {
            (a.get("email") or "").lower() for a in attendees
        }
        nuevos = []
        for em in emails:
            if em and em.lower() not in existentes:
                attendees.append({"email": em})
                existentes.add(em.lower())
                nuevos.append(em)
        ev2 = (
            cal.events()
            .patch(
                calendarId=calendar_id,
                eventId=event_id,
                body={"attendees": attendees},
                sendUpdates="all",
            )
            .execute()
        )
        _audit(
            "calendar_anadir_asistentes",
            f"calendar={calendar_id} event={event_id} nuevos={nuevos}",
        )
        return ev2.get("attendees", []) or []
    except Exception as e:
        return _err("calendar_anadir_asistentes", e)



# =========================================================
# BLOQUE: Slides + Tasks + Contactos (Google Workspace)
# Se APPENDEA a server.py (FastMCP). No redefine nada existente.
# =========================================================


def _get_slides():
    """Devuelve (y cachea) el servicio de Google Slides v1."""
    if not hasattr(_get_slides, "_svc"):
        _get_slides._svc = build(
            "slides", "v1", credentials=_build_creds(), cache_discovery=False
        )
    return _get_slides._svc


def _get_tasks():
    """Devuelve (y cachea) el servicio de Google Tasks v1."""
    if not hasattr(_get_tasks, "_svc"):
        _get_tasks._svc = build(
            "tasks", "v1", credentials=_build_creds(), cache_discovery=False
        )
    return _get_tasks._svc


def _get_people():
    """Devuelve (y cachea) el servicio de Google People v1."""
    if not hasattr(_get_people, "_svc"):
        _get_people._svc = build(
            "people", "v1", credentials=_build_creds(), cache_discovery=False
        )
    return _get_people._svc


# ------------------------- SLIDES -------------------------

@mcp.tool()
def slides_crear(titulo: str, diapositivas: Optional[list] = None) -> dict:
    """Crea una presentación de Google Slides.

    Args:
        titulo: Título de la presentación.
        diapositivas: Lista opcional de dicts {'titulo':.., 'cuerpo':..};
            por cada elemento se añade una diapositiva con layout
            TITLE_AND_BODY rellenando ambos placeholders.

    Devuelve el presentationId y la url de edición.
    """
    try:
        svc = _get_slides()
        pres = svc.presentations().create(body={"title": titulo}).execute()
        pid = pres.get("presentationId")
        if diapositivas:
            requests = []
            for idx, dia in enumerate(diapositivas):
                slide_id = "slide_%d" % idx
                title_id = "title_%d" % idx
                body_id = "body_%d" % idx
                requests.append({
                    "createSlide": {
                        "objectId": slide_id,
                        "slideLayoutReference": {
                            "predefinedLayout": "TITLE_AND_BODY"
                        },
                        "placeholderIdMappings": [
                            {
                                "layoutPlaceholder": {"type": "TITLE", "index": 0},
                                "objectId": title_id,
                            },
                            {
                                "layoutPlaceholder": {"type": "BODY", "index": 0},
                                "objectId": body_id,
                            },
                        ],
                    }
                })
                if dia.get("titulo"):
                    requests.append({
                        "insertText": {
                            "objectId": title_id,
                            "text": str(dia.get("titulo")),
                        }
                    })
                if dia.get("cuerpo"):
                    requests.append({
                        "insertText": {
                            "objectId": body_id,
                            "text": str(dia.get("cuerpo")),
                        }
                    })
            if requests:
                svc.presentations().batchUpdate(
                    presentationId=pid, body={"requests": requests}
                ).execute()
        url = "https://docs.google.com/presentation/d/%s/edit" % pid
        _audit("slides_crear", "%s - %s" % (pid, titulo))
        return {"presentationId": pid, "url": url}
    except Exception as e:
        return _err("slides_crear", e)


@mcp.tool()
def slides_leer(presentation_id: str, max_chars: int = 20000) -> dict:
    """Extrae el texto de una presentación de Google Slides.

    Recorre todas las diapositivas y sus elementos de forma (shapes),
    concatenando el contenido de shape.text.textElements[].textRun.content.

    Args:
        presentation_id: ID de la presentación.
        max_chars: Máximo de caracteres a devolver.

    Devuelve el texto extraído.
    """
    try:
        svc = _get_slides()
        pres = svc.presentations().get(presentationId=presentation_id).execute()
        partes = []
        for slide in pres.get("slides", []):
            for elem in slide.get("pageElements", []):
                shape = elem.get("shape")
                if not shape:
                    continue
                text = shape.get("text")
                if not text:
                    continue
                for te in text.get("textElements", []):
                    tr = te.get("textRun")
                    if tr and tr.get("content"):
                        partes.append(tr["content"])
        texto = "".join(partes)[:max_chars]
        return {"presentationId": presentation_id, "texto": texto}
    except Exception as e:
        return _err("slides_leer", e)


@mcp.tool()
def slides_anadir_diapositiva(
    presentation_id: str, titulo: str, cuerpo: str = ""
) -> dict:
    """Añade una diapositiva (layout TITLE_AND_BODY) a una presentación.

    Args:
        presentation_id: ID de la presentación.
        titulo: Texto para el placeholder de título.
        cuerpo: Texto opcional para el placeholder de cuerpo.

    Devuelve ok.
    """
    try:
        svc = _get_slides()
        suf = base64.urlsafe_b64encode(os.urandom(9)).decode().rstrip("=")
        slide_id = "slide_" + suf
        title_id = "title_" + suf
        body_id = "body_" + suf
        requests = [
            {
                "createSlide": {
                    "objectId": slide_id,
                    "slideLayoutReference": {
                        "predefinedLayout": "TITLE_AND_BODY"
                    },
                    "placeholderIdMappings": [
                        {
                            "layoutPlaceholder": {"type": "TITLE", "index": 0},
                            "objectId": title_id,
                        },
                        {
                            "layoutPlaceholder": {"type": "BODY", "index": 0},
                            "objectId": body_id,
                        },
                    ],
                }
            },
            {"insertText": {"objectId": title_id, "text": titulo}},
        ]
        if cuerpo:
            requests.append(
                {"insertText": {"objectId": body_id, "text": cuerpo}}
            )
        svc.presentations().batchUpdate(
            presentationId=presentation_id, body={"requests": requests}
        ).execute()
        _audit("slides_anadir_diapositiva", "%s - %s" % (presentation_id, titulo))
        return {"ok": True}
    except Exception as e:
        return _err("slides_anadir_diapositiva", e)


# ------------------------- TASKS --------------------------

@mcp.tool()
def tasks_listar(max_results: int = 50) -> dict:
    """Lista las tareas de la lista por defecto (@default).

    Args:
        max_results: Número máximo de tareas a devolver.

    Devuelve id, título, estado y vencimiento de cada tarea.
    """
    try:
        svc = _get_tasks()
        res = svc.tasks().list(
            tasklist="@default", maxResults=max_results
        ).execute()
        items = []
        for t in res.get("items", []):
            items.append({
                "id": t.get("id"),
                "titulo": t.get("title"),
                "estado": t.get("status"),
                "vencimiento": t.get("due"),
            })
        return {"tareas": items}
    except Exception as e:
        return _err("tasks_listar", e)


@mcp.tool()
def tasks_crear(titulo: str, fecha: Optional[str] = None, notas: str = "") -> dict:
    """Crea una tarea en la lista por defecto (@default).

    Args:
        titulo: Título de la tarea.
        fecha: Vencimiento en RFC3339; si llega solo 'AAAA-MM-DD' se
            convierte a 'AAAA-MM-DDT00:00:00.000Z'.
        notas: Notas opcionales.

    Devuelve el id de la tarea creada.
    """
    try:
        svc = _get_tasks()
        due = fecha
        if fecha and len(fecha) == 10 and "T" not in fecha:
            due = fecha + "T00:00:00.000Z"
        body = {"title": titulo, "notes": notas}
        if due:
            body["due"] = due
        t = svc.tasks().insert(tasklist="@default", body=body).execute()
        _audit("tasks_crear", "%s - %s" % (t.get("id"), titulo))
        return {"id": t.get("id")}
    except Exception as e:
        return _err("tasks_crear", e)


@mcp.tool()
def tasks_completar(task_id: str) -> dict:
    """Marca una tarea como completada (status='completed').

    Args:
        task_id: ID de la tarea.

    Devuelve ok.
    """
    try:
        svc = _get_tasks()
        svc.tasks().patch(
            tasklist="@default", task=task_id, body={"status": "completed"}
        ).execute()
        _audit("tasks_completar", task_id)
        return {"ok": True}
    except Exception as e:
        return _err("tasks_completar", e)


# ----------------------- CONTACTOS ------------------------

def _people_a_dict(p: dict) -> dict:
    """Normaliza una persona de People API a nombre/emails/teléfonos."""
    nombre = ""
    if p.get("names"):
        nombre = p["names"][0].get("displayName", "")
    emails = [e.get("value") for e in p.get("emailAddresses", []) if e.get("value")]
    telefonos = [t.get("value") for t in p.get("phoneNumbers", []) if t.get("value")]
    return {"nombre": nombre, "emails": emails, "telefonos": telefonos}


@mcp.tool()
def contactos_buscar(texto: str, max_results: int = 10) -> dict:
    """Busca contactos por texto (People API searchContacts).

    Args:
        texto: Consulta de búsqueda.
        max_results: Número máximo de resultados.

    Devuelve nombre, emails y teléfonos de cada contacto.
    """
    try:
        svc = _get_people()
        res = svc.people().searchContacts(
            query=texto,
            readMask="names,emailAddresses,phoneNumbers",
            pageSize=max_results,
        ).execute()
        contactos = [_people_a_dict(r.get("person", {})) for r in res.get("results", [])]
        return {"contactos": contactos}
    except Exception as e:
        return _err("contactos_buscar", e)


@mcp.tool()
def contactos_listar(max_results: int = 50) -> dict:
    """Lista los contactos del usuario (People API connections.list).

    Args:
        max_results: Número máximo de contactos.

    Devuelve la lista con nombre, emails y teléfonos.
    """
    try:
        svc = _get_people()
        res = svc.people().connections().list(
            resourceName="people/me",
            personFields="names,emailAddresses,phoneNumbers",
            pageSize=max_results,
        ).execute()
        contactos = [_people_a_dict(p) for p in res.get("connections", [])]
        return {"contactos": contactos}
    except Exception as e:
        return _err("contactos_listar", e)




# ============================================================
# BLOQUE: Unidades compartidas + PDF legal
# ============================================================


@mcp.tool()
def drive_unidades_compartidas() -> dict:
    """Lista las unidades compartidas (shared drives) accesibles.

    Devuelve el id y el nombre de cada unidad compartida.
    """
    try:
        service = _get_service()
        resp = service.drives().list(
            pageSize=100,
            fields="drives(id,name)",
        ).execute()
        unidades = [
            {"id": d.get("id"), "nombre": d.get("name")}
            for d in resp.get("drives", [])
        ]
        return {"ok": True, "total": len(unidades), "unidades": unidades}
    except Exception as e:
        return _err("drive_unidades_compartidas", e)


def _pdf_descargar_bytes(file_id: str) -> "io.BytesIO":
    """Descarga un archivo de Drive a un io.BytesIO (uso interno)."""
    service = _get_service()
    buffer = io.BytesIO()
    request = service.files().get_media(fileId=file_id)
    downloader = MediaIoBaseDownload(buffer, request)
    done = False
    while not done:
        _, done = downloader.next_chunk()
    buffer.seek(0)
    return buffer


def _pdf_subir_bytes(buffer: "io.BytesIO", nombre_salida: str, parent_id: str) -> str:
    """Sube un io.BytesIO como PDF a Drive y devuelve el id (uso interno)."""
    service = _get_service()
    buffer.seek(0)
    metadata = {"name": nombre_salida}
    if parent_id:
        metadata["parents"] = [parent_id]
    media = MediaIoBaseUpload(buffer, mimetype="application/pdf", resumable=True)
    creado = service.files().create(
        body=metadata,
        media_body=media,
        fields=FILE_FIELDS,
        supportsAllDrives=True,
    ).execute()
    return creado.get("id")


@mcp.tool()
def drive_pdf_unir(file_ids: list, nombre_salida: str, parent_id: str = "root") -> dict:
    """Une varios PDF de Drive en uno solo y lo sube.

    Descarga cada PDF indicado en 'file_ids', los concatena con pypdf
    en el orden dado y sube el PDF resultante a 'parent_id'.
    Devuelve el id del PDF creado.
    """
    try:
        from pypdf import PdfWriter

        if not file_ids:
            return _err("drive_pdf_unir", ValueError("file_ids vacio"))

        writer = PdfWriter()
        for fid in file_ids:
            origen = _pdf_descargar_bytes(fid)
            writer.append(origen)

        salida = io.BytesIO()
        writer.write(salida)
        writer.close()

        nuevo_id = _pdf_subir_bytes(salida, nombre_salida, parent_id)
        _audit(
            "drive_pdf_unir",
            f"unidos {len(file_ids)} PDF en '{nombre_salida}' (id={nuevo_id})",
        )
        return {"ok": True, "id": nuevo_id, "nombre": nombre_salida}
    except Exception as e:
        return _err("drive_pdf_unir", e)


@mcp.tool()
def drive_pdf_marca_agua(
    file_id: str,
    texto: str,
    nombre_salida: str,
    parent_id: str = "root",
) -> dict:
    """Aplica una marca de agua de texto en diagonal a un PDF y lo sube.

    Genera con reportlab una pagina con 'texto' en diagonal, en gris
    translucido y centrado, y la superpone sobre CADA pagina del PDF
    original. Sube el resultado a 'parent_id'. Devuelve el id creado.
    """
    try:
        from pypdf import PdfReader, PdfWriter
        from reportlab.pdfgen import canvas
        from reportlab.lib.colors import Color

        origen = _pdf_descargar_bytes(file_id)
        lector = PdfReader(origen)
        escritor = PdfWriter()

        for pagina in lector.pages:
            ancho = float(pagina.mediabox.width)
            alto = float(pagina.mediabox.height)

            overlay_buf = io.BytesIO()
            c = canvas.Canvas(overlay_buf, pagesize=(ancho, alto))
            c.saveState()
            c.setFont("Helvetica-Bold", 60)
            c.setFillColor(Color(0.5, 0.5, 0.5, alpha=0.3))
            c.translate(ancho / 2.0, alto / 2.0)
            c.rotate(45)
            c.drawCentredString(0, 0, texto)
            c.restoreState()
            c.save()
            overlay_buf.seek(0)

            overlay_pagina = PdfReader(overlay_buf).pages[0]
            pagina.merge_page(overlay_pagina)
            escritor.add_page(pagina)

        salida = io.BytesIO()
        escritor.write(salida)
        escritor.close()

        nuevo_id = _pdf_subir_bytes(salida, nombre_salida, parent_id)
        _audit(
            "drive_pdf_marca_agua",
            f"marca de agua '{texto}' sobre {file_id} -> '{nombre_salida}' (id={nuevo_id})",
        )
        return {"ok": True, "id": nuevo_id, "nombre": nombre_salida}
    except Exception as e:
        return _err("drive_pdf_marca_agua", e)


@mcp.tool()
def drive_pdf_foliar(
    file_id: str,
    nombre_salida: str,
    parent_id: str = "root",
    prefijo: str = "",
) -> dict:
    """Anade foliado tipo Bates (numeracion de paginas) a un PDF y lo sube.

    Superpone en cada pagina, abajo a la derecha, un sello con el formato
    f"{prefijo}{n:04d}" mediante un overlay de reportlab del tamano de la
    pagina. Sube el resultado a 'parent_id'. Devuelve el id creado.
    """
    try:
        from pypdf import PdfReader, PdfWriter
        from reportlab.pdfgen import canvas

        origen = _pdf_descargar_bytes(file_id)
        lector = PdfReader(origen)
        escritor = PdfWriter()

        for indice, pagina in enumerate(lector.pages, start=1):
            ancho = float(pagina.mediabox.width)
            alto = float(pagina.mediabox.height)
            sello = f"{prefijo}{indice:04d}"

            overlay_buf = io.BytesIO()
            c = canvas.Canvas(overlay_buf, pagesize=(ancho, alto))
            c.setFont("Helvetica", 9)
            c.setFillColorRGB(0, 0, 0)
            c.drawRightString(ancho - 36, 24, sello)
            c.save()
            overlay_buf.seek(0)

            overlay_pagina = PdfReader(overlay_buf).pages[0]
            pagina.merge_page(overlay_pagina)
            escritor.add_page(pagina)

        salida = io.BytesIO()
        escritor.write(salida)
        escritor.close()

        nuevo_id = _pdf_subir_bytes(salida, nombre_salida, parent_id)
        _audit(
            "drive_pdf_foliar",
            f"foliado '{prefijo}' sobre {file_id} -> '{nombre_salida}' (id={nuevo_id})",
        )
        return {"ok": True, "id": nuevo_id, "nombre": nombre_salida}
    except Exception as e:
        return _err("drive_pdf_foliar", e)


# ============================================================================
# BLOQUE: Cálculo de PLAZOS PROCESALES (jurisdicción española)
# ----------------------------------------------------------------------------
# Python puro con datetime. Se APPENDEA a server.py (FastMCP).
# Requiere ya definidos: mcp, Optional (typing) y el helper _err(action, e).
# NOTA: los festivos AUTONÓMICOS y LOCALES no se incluyen aquí; pueden
#       pasarse mediante el parámetro 'festivos_extra' (lista de 'AAAA-MM-DD').
# ADVERTENCIA: herramienta de ayuda; no sustituye la verificación por un
#              profesional del cómputo aplicable a cada procedimiento.
# ============================================================================

# Festivos NACIONALES de España para 2025 y 2026 (fechas fijas + Viernes Santo).
FESTIVOS_NACIONALES_ES = [
    # 2025
    "2025-01-01",  # Año Nuevo
    "2025-01-06",  # Epifanía / Reyes
    "2025-04-18",  # Viernes Santo
    "2025-05-01",  # Fiesta del Trabajo
    "2025-08-15",  # Asunción de la Virgen
    "2025-10-12",  # Fiesta Nacional de España
    "2025-11-01",  # Todos los Santos
    "2025-12-06",  # Día de la Constitución
    "2025-12-08",  # Inmaculada Concepción
    "2025-12-25",  # Natividad del Señor
    # 2026
    "2026-01-01",  # Año Nuevo
    "2026-01-06",  # Epifanía / Reyes
    "2026-04-03",  # Viernes Santo
    "2026-05-01",  # Fiesta del Trabajo
    "2026-08-15",  # Asunción de la Virgen
    "2026-10-12",  # Fiesta Nacional de España
    "2026-11-01",  # Todos los Santos
    "2026-12-06",  # Día de la Constitución
    "2026-12-08",  # Inmaculada Concepción
    "2026-12-25",  # Natividad del Señor
]

# Nombres de los días de la semana en español (0 = lunes ... 6 = domingo).
_DIAS_SEMANA_ES = [
    "lunes", "martes", "miércoles", "jueves", "viernes", "sábado", "domingo",
]


def _es_dia_inhabil(d, sabado_inhabil, agosto_inhabil, festivos):
    """Devuelve True si la fecha 'd' (datetime.date) es día inhábil."""
    wd = d.weekday()  # lunes=0 ... domingo=6
    if wd == 6:  # el domingo es SIEMPRE inhábil
        return True
    if sabado_inhabil and wd == 5:  # sábado inhábil si procede
        return True
    if agosto_inhabil and d.month == 8:  # todo agosto (plazos procesales)
        return True
    if d.isoformat() in festivos:  # festivos nacionales + extra
        return True
    return False


def _construir_festivos(festivos_extra):
    """Combina los festivos nacionales con los extra recibidos (o None)."""
    festivos = set(FESTIVOS_NACIONALES_ES)
    if festivos_extra:
        festivos.update(str(f) for f in festivos_extra)
    return festivos


@mcp.tool()
def plazo_dias_habiles(
    fecha_inicio: str,
    dias: int,
    sabado_inhabil: bool = True,
    agosto_inhabil: bool = False,
    festivos_extra: Optional[list] = None,
):
    """Suma días HÁBILES a una fecha para obtener el vencimiento del plazo.

    Partiendo de 'fecha_inicio' (AAAA-MM-DD) suma 'dias' días hábiles: salta
    siempre los domingos; los sábados si 'sabado_inhabil'; los festivos
    nacionales más 'festivos_extra'; y todo el mes de agosto si
    'agosto_inhabil' (propio de plazos procesales). Devuelve la fecha de
    vencimiento (AAAA-MM-DD), el día de la semana y los días naturales
    transcurridos.
    """
    try:
        import datetime
        festivos = _construir_festivos(festivos_extra)
        inicio = datetime.date.fromisoformat(fecha_inicio)
        actual = inicio
        contados = 0
        while contados < dias:
            actual = actual + datetime.timedelta(days=1)
            if not _es_dia_inhabil(actual, sabado_inhabil, agosto_inhabil, festivos):
                contados += 1
        return {
            "fecha_inicio": inicio.isoformat(),
            "dias_habiles_sumados": dias,
            "fecha_vencimiento": actual.isoformat(),
            "dia_semana": _DIAS_SEMANA_ES[actual.weekday()],
            "dias_naturales_transcurridos": (actual - inicio).days,
            "sabado_inhabil": sabado_inhabil,
            "agosto_inhabil": agosto_inhabil,
        }
    except Exception as e:
        return _err("plazo_dias_habiles", e)


@mcp.tool()
def plazo_vencimiento_notificacion(
    fecha_notificacion: str,
    dias_habiles: int,
    agosto_inhabil: bool = False,
    festivos_extra: Optional[list] = None,
):
    """Calcula el vencimiento de un plazo desde la notificación.

    Atajo que aplica la regla habitual: el cómputo empieza el día siguiente
    al de la notificación (AAAA-MM-DD). Los sábados se consideran inhábiles.
    Devuelve la fecha de vencimiento, el día de la semana y una nota
    recordando que es una ayuda y no sustituye la verificación del
    profesional para cada procedimiento concreto.
    """
    try:
        import datetime
        festivos = _construir_festivos(festivos_extra)
        inicio = datetime.date.fromisoformat(fecha_notificacion)
        actual = inicio
        contados = 0
        while contados < dias_habiles:
            actual = actual + datetime.timedelta(days=1)
            if not _es_dia_inhabil(actual, True, agosto_inhabil, festivos):
                contados += 1
        return {
            "fecha_notificacion": inicio.isoformat(),
            "dias_habiles": dias_habiles,
            "fecha_vencimiento": actual.isoformat(),
            "dia_semana": _DIAS_SEMANA_ES[actual.weekday()],
            "nota": (
                "El cómputo se inicia el día siguiente al de la notificación. "
                "Herramienta de ayuda: no sustituye la verificación del cómputo "
                "por el profesional según el procedimiento aplicable."
            ),
        }
    except Exception as e:
        return _err("plazo_vencimiento_notificacion", e)


@mcp.tool()
def dias_habiles_entre(
    fecha_a: str,
    fecha_b: str,
    sabado_inhabil: bool = True,
    agosto_inhabil: bool = False,
    festivos_extra: Optional[list] = None,
):
    """Cuenta los días HÁBILES entre dos fechas (AAAA-MM-DD).

    Cuenta los días hábiles del intervalo abierto por la izquierda y cerrado
    por la derecha (desde el día siguiente a la fecha menor hasta la fecha
    mayor, ambos según los criterios de inhabilidad). Aplica los mismos
    criterios que las demás herramientas (domingos, sábados si procede,
    festivos nacionales + extra y agosto si procede). El orden de las fechas
    es indiferente.
    """
    try:
        import datetime
        festivos = _construir_festivos(festivos_extra)
        d1 = datetime.date.fromisoformat(fecha_a)
        d2 = datetime.date.fromisoformat(fecha_b)
        inicio, fin = (d1, d2) if d1 <= d2 else (d2, d1)
        contados = 0
        actual = inicio
        while actual < fin:
            actual = actual + datetime.timedelta(days=1)
            if not _es_dia_inhabil(actual, sabado_inhabil, agosto_inhabil, festivos):
                contados += 1
        return {
            "fecha_a": inicio.isoformat(),
            "fecha_b": fin.isoformat(),
            "dias_habiles": contados,
            "dias_naturales": (fin - inicio).days,
            "sabado_inhabil": sabado_inhabil,
            "agosto_inhabil": agosto_inhabil,
        }
    except Exception as e:
        return _err("dias_habiles_entre", e)



# --------------------------------------------------------------------------- #
# BLOQUE PREMIUM: Doc-merge por lotes + comandos combinados (360, adjuntos, Meet)
# --------------------------------------------------------------------------- #

_DOC_MIME = "application/vnd.google-apps.document"


@mcp.tool()
def drive_generar_desde_plantilla_lote(template_id: str, filas: list,
                                       parent_id: str = "root",
                                       patron_nombre: str = "Documento") -> dict:
    """Genera un Documento de Google por cada fila a partir de una plantilla (mail-merge).

    Para CADA dict de `filas`: copia el Documento `template_id`, nombra la copia con
    `patron_nombre` formateado con la fila (p.ej. 'Contrato {nombre}' -> patron_nombre.format(**fila)),
    reemplaza en la copia cada marcador '{{clave}}' por su valor mediante replaceAllText,
    y mueve la copia a `parent_id`. Devuelve la lista de documentos creados (id, nombre).

    `filas` es una lista de dicts, p.ej.:
        [{"nombre": "Ana", "puesto": "Analista"}, {"nombre": "Luis", "puesto": "Jefe"}]
    """
    try:
        drv = _get_service()
        docs = _get_docs()
        creados = []
        for fila in filas:
            try:
                nombre = patron_nombre.format(**fila)
            except Exception:
                nombre = patron_nombre
            copia = drv.files().copy(
                fileId=template_id,
                body={"name": nombre},
                fields="id, name, parents",
                supportsAllDrives=True,
            ).execute()
            copia_id = copia["id"]

            requests = []
            for clave, valor in fila.items():
                requests.append({
                    "replaceAllText": {
                        "containsText": {"text": "{{" + str(clave) + "}}", "matchCase": True},
                        "replaceText": str(valor),
                    }
                })
            if requests:
                docs.documents().batchUpdate(
                    documentId=copia_id,
                    body={"requests": requests},
                ).execute()

            prev = ",".join(copia.get("parents", []))
            drv.files().update(
                fileId=copia_id,
                addParents=parent_id,
                removeParents=prev,
                fields="id, parents",
                supportsAllDrives=True,
            ).execute()

            _audit("drive_generar_desde_plantilla_lote",
                   {"template_id": template_id, "doc_id": copia_id, "nombre": nombre,
                    "parent_id": parent_id})
            creados.append({"id": copia_id, "nombre": nombre})
        return {"ok": True, "generados": len(creados), "documentos": creados}
    except Exception as e:
        return _err("drive_generar_desde_plantilla_lote", e)


@mcp.tool()
def cliente_360(nombre_cliente: str, max_items: int = 10) -> dict:
    """Vista 360 de un cliente: reune en un solo golpe su rastro en Drive, Gmail y Calendar.

    Devuelve un dict con tres listas resumidas:
      - `drive`: archivos/carpetas cuyo nombre contiene `nombre_cliente`.
      - `correos`: correos recientes de/para el cliente (From, Subject, Date).
      - `eventos`: proximos eventos de calendario que lo mencionan.
    """
    try:
        tope = max(1, min(max_items, 50))
        resultado = {"cliente": nombre_cliente, "drive": [], "correos": [], "eventos": []}

        # (a) Drive: nombre contiene el cliente
        try:
            drv = _get_service()
            safe = nombre_cliente.replace("'", "\\'")
            q = "name contains '{}' and trashed = false".format(safe)
            res = drv.files().list(q=q, spaces="drive", fields="files({})".format(FILE_FIELDS),
                                   pageSize=tope, supportsAllDrives=True,
                                   includeItemsFromAllDrives=True).execute()
            resultado["drive"] = [_file(f) for f in res.get("files", [])]
        except Exception as e:
            resultado["drive_error"] = str(e)

        # (b) Gmail: correos donde aparezca el cliente
        try:
            gm = _get_gmail()
            gq = 'from:{0} OR to:{0} OR "{0}"'.format(nombre_cliente)
            lst = gm.users().messages().list(userId="me", q=gq, maxResults=tope).execute()
            for m in lst.get("messages", []):
                full = gm.users().messages().get(
                    userId="me", id=m["id"], format="metadata",
                    metadataHeaders=["From", "To", "Subject", "Date"]).execute()
                h = {x["name"]: x["value"] for x in full.get("payload", {}).get("headers", [])}
                resultado["correos"].append({
                    "id": m["id"], "from": h.get("From"), "to": h.get("To"),
                    "subject": h.get("Subject"), "date": h.get("Date"),
                    "snippet": full.get("snippet", "")[:160]})
        except Exception as e:
            resultado["correos_error"] = str(e)

        # (c) Calendar: proximos eventos que lo mencionan
        try:
            import datetime as _dt
            cal = _get_cal()
            ahora = _dt.datetime.utcnow().strftime("%Y-%m-%dT%H:%M:%SZ")
            ev = cal.events().list(calendarId="primary", q=nombre_cliente, timeMin=ahora,
                                   singleEvents=True, orderBy="startTime",
                                   maxResults=tope).execute()
            for e in ev.get("items", []):
                st = e.get("start") or {}
                resultado["eventos"].append({
                    "id": e.get("id"), "summary": e.get("summary"),
                    "inicio": st.get("dateTime") or st.get("date"),
                    "link": e.get("htmlLink")})
        except Exception as e:
            resultado["eventos_error"] = str(e)

        resultado["ok"] = True
        return resultado
    except Exception as e:
        return _err("cliente_360", e)


@mcp.tool()
def archivar_adjuntos_recientes(dias: int = 7, parent_id: str = "root") -> dict:
    """Archiva en Drive todos los adjuntos de los correos recientes con adjunto.

    Busca en Gmail 'has:attachment newer_than:{dias}d', descarga cada adjunto y lo sube
    a la carpeta `parent_id` de Drive. Devuelve la lista de adjuntos guardados (nombre, drive_id).
    """
    try:
        gm = _get_gmail()
        drv = _get_service()

        def _walk(payload):
            pila = [payload]
            while pila:
                p = pila.pop()
                for c in (p.get("parts") or []):
                    pila.append(c)
                yield p

        query = "has:attachment newer_than:{}d".format(max(1, dias))
        lst = gm.users().messages().list(userId="me", q=query, maxResults=100).execute()
        guardados = []
        for m in lst.get("messages", []):
            full = gm.users().messages().get(userId="me", id=m["id"], format="full").execute()
            for p in _walk(full.get("payload", {})):
                fn = p.get("filename")
                body = p.get("body", {})
                if fn and body.get("attachmentId"):
                    att = gm.users().messages().attachments().get(
                        userId="me", messageId=m["id"], id=body["attachmentId"]).execute()
                    data = base64.urlsafe_b64decode(att["data"])
                    media = MediaIoBaseUpload(io.BytesIO(data),
                                              mimetype="application/octet-stream",
                                              resumable=False)
                    r = drv.files().create(body={"name": fn, "parents": [parent_id]},
                                           media_body=media, fields="id, name",
                                           supportsAllDrives=True).execute()
                    _audit("archivar_adjuntos_recientes",
                           {"message_id": m["id"], "file": fn, "drive_id": r["id"],
                            "parent_id": parent_id})
                    guardados.append({"nombre": fn, "drive_id": r["id"]})
        return {"ok": True, "guardados": len(guardados), "adjuntos": guardados}
    except Exception as e:
        return _err("archivar_adjuntos_recientes", e)


@mcp.tool()
def agendar_reunion_cliente(summary: str, start: str, end: str, emails: list,
                            descripcion: str = "", calendar_id: str = "primary") -> dict:
    """Agenda una reunion con enlace de Google Meet e invita a los asistentes.

    Crea un evento con videollamada de Google Meet (conferenceData/hangoutsMeet) e invita
    a `emails` enviando la invitacion (sendUpdates='all'). start/end en RFC3339 con zona
    (2026-08-05T10:00:00+01:00). Devuelve el id del evento y el enlace de Meet.
    """
    try:
        import uuid as _uuid
        cal = _get_cal()
        body = {
            "summary": summary,
            "description": descripcion,
            "start": {"dateTime": start},
            "end": {"dateTime": end},
            "attendees": [{"email": a} for a in (emails or [])],
            "conferenceData": {
                "createRequest": {
                    "requestId": str(_uuid.uuid4()),
                    "conferenceSolutionKey": {"type": "hangoutsMeet"},
                }
            },
        }
        e = cal.events().insert(calendarId=calendar_id, body=body,
                                conferenceDataVersion=1, sendUpdates="all").execute()
        _audit("agendar_reunion_cliente",
               {"summary": summary, "start": start, "event_id": e.get("id"),
                "emails": list(emails or [])})
        return {"ok": True, "id": e.get("id"), "meet": e.get("hangoutLink"),
                "link": e.get("htmlLink")}
    except Exception as e:
        return _err("agendar_reunion_cliente", e)



# --------------------------------------------------------------------------- #
# GMAIL extra: hilos, estado leido/no leido y carpetas (etiquetas)
# --------------------------------------------------------------------------- #

@mcp.tool()
def gmail_leer_hilo(thread_id: str, max_chars: int = 30000):
    """Lee un HILO completo de Gmail (la conversacion con todas sus respuestas).
    Devuelve, por cada mensaje, remitente, fecha, asunto y texto. Acepta el id del
    hilo o el id de cualquier mensaje del hilo."""
    try:
        import base64 as _b64
        svc = _get_gmail()
        try:
            th = svc.users().threads().get(userId='me', id=thread_id, format='full').execute()
        except Exception:
            m = svc.users().messages().get(userId='me', id=thread_id, format='minimal').execute()
            th = svc.users().threads().get(userId='me', id=m['threadId'], format='full').execute()
        msgs = []
        for m in th.get('messages', []):
            payload = m.get('payload', {})
            h = {x['name']: x['value'] for x in payload.get('headers', [])}
            text = ''
            stack = [payload]
            while stack:
                p = stack.pop()
                for c in (p.get('parts') or []):
                    stack.append(c)
                b = p.get('body', {})
                if p.get('mimeType') == 'text/plain' and b.get('data'):
                    text += _b64.urlsafe_b64decode(b['data']).decode('utf-8', 'replace')
            msgs.append({'id': m.get('id'), 'from': h.get('From'), 'date': h.get('Date'),
                         'subject': h.get('Subject'), 'texto': text[:max_chars]})
        return {"ok": True, "thread_id": th.get('id'), "mensajes": len(msgs), "conversacion": msgs}
    except Exception as e:
        return _err("gmail_leer_hilo", e)


@mcp.tool()
def gmail_marcar_leido(message_id: str, leido: bool = True):
    """Marca un correo como LEIDO (leido=True) o como NO LEIDO (leido=False)."""
    try:
        svc = _get_gmail()
        body = {"removeLabelIds": ["UNREAD"]} if leido else {"addLabelIds": ["UNREAD"]}
        svc.users().messages().modify(userId='me', id=message_id, body=body).execute()
        _audit("gmail_marcar_leido", {"message_id": message_id, "leido": leido})
        return {"ok": True, "message_id": message_id, "leido": leido}
    except Exception as e:
        return _err("gmail_marcar_leido", e)


@mcp.tool()
def gmail_listar_carpetas():
    """Lista las carpetas/etiquetas de Gmail (nombre e id)."""
    try:
        svc = _get_gmail()
        res = svc.users().labels().list(userId='me').execute()
        labels = [{"id": l.get("id"), "nombre": l.get("name"), "tipo": l.get("type")}
                  for l in res.get("labels", [])]
        return {"ok": True, "count": len(labels), "carpetas": labels}
    except Exception as e:
        return _err("gmail_listar_carpetas", e)


@mcp.tool()
def gmail_crear_carpeta(nombre: str):
    """Crea una carpeta/etiqueta nueva en Gmail. Devuelve su id."""
    try:
        svc = _get_gmail()
        l = svc.users().labels().create(userId='me', body={"name": nombre,
            "labelListVisibility": "labelShow", "messageListVisibility": "show"}).execute()
        _audit("gmail_crear_carpeta", {"nombre": nombre, "id": l.get("id")})
        return {"ok": True, "id": l.get("id"), "nombre": nombre}
    except Exception as e:
        return _err("gmail_crear_carpeta", e)


@mcp.tool()
def gmail_archivar_en(message_id: str, carpeta: str, crear_si_no_existe: bool = True,
                      sacar_de_bandeja: bool = False):
    """Guarda/mueve un correo a una carpeta (etiqueta) por su NOMBRE. Si no existe y
    crear_si_no_existe=True, la crea. Con sacar_de_bandeja=True lo quita de la Bandeja de
    entrada (archivar de verdad)."""
    try:
        svc = _get_gmail()
        res = svc.users().labels().list(userId='me').execute()
        lid = None
        for l in res.get("labels", []):
            if (l.get("name") or "").lower() == carpeta.lower():
                lid = l.get("id"); break
        if not lid:
            if not crear_si_no_existe:
                return {"ok": False, "action": "gmail_archivar_en",
                        "error": "La carpeta no existe y crear_si_no_existe=False."}
            l = svc.users().labels().create(userId='me', body={"name": carpeta,
                "labelListVisibility": "labelShow", "messageListVisibility": "show"}).execute()
            lid = l.get("id")
        body = {"addLabelIds": [lid]}
        if sacar_de_bandeja:
            body["removeLabelIds"] = ["INBOX"]
        svc.users().messages().modify(userId='me', id=message_id, body=body).execute()
        _audit("gmail_archivar_en", {"message_id": message_id, "carpeta": carpeta,
               "sacar_de_bandeja": sacar_de_bandeja})
        return {"ok": True, "message_id": message_id, "carpeta": carpeta, "label_id": lid}
    except Exception as e:
        return _err("gmail_archivar_en", e)


# =====================================================================
# BLOQUE: Gmail acciones basicas (Gmail v1)
# Se APPENDEA a server.py. No anade imports de nivel modulo ni redefine nada.
# Usa helpers ya existentes: _get_gmail(), _err(), _audit().
# =====================================================================


@mcp.tool()
def gmail_responder_a_todos(message_id: str, texto: str):
    """Responde a TODOS los participantes de un correo dentro de su MISMO hilo.

    Obtiene el mensaje original (From, To, Cc, Subject, Message-Id, References,
    threadId), compone la respuesta con destinatarios el remitente original mas
    todos los To y Cc originales EXCLUYENDO la propia direccion del usuario,
    asunto con prefijo 'Re:' y cabeceras In-Reply-To y References, y la envia en
    el mismo threadId. Devuelve el id del mensaje enviado.
    """
    try:
        from email.mime.text import MIMEText
        from email.utils import getaddresses, formataddr
        svc = _get_gmail()
        yo = svc.users().getProfile(userId='me').execute().get('emailAddress', '') or ''
        original = svc.users().messages().get(
            userId='me', id=message_id, format='metadata',
            metadataHeaders=['From', 'To', 'Cc', 'Subject', 'Message-Id', 'References']).execute()
        thread_id = original.get('threadId')
        headers = original.get('payload', {}).get('headers', [])

        def _h(name):
            for h in headers:
                if h.get('name', '').lower() == name.lower():
                    return h.get('value')
            return None

        remitente = _h('From') or ''
        to_orig = _h('To') or ''
        cc_orig = _h('Cc') or ''
        asunto_orig = _h('Subject') or ''
        msg_id_hdr = _h('Message-Id')
        refs = _h('References')

        destinatarios = []
        vistos = set()
        yo_l = yo.strip().lower()
        for nombre, correo in getaddresses([remitente, to_orig, cc_orig]):
            correo_l = (correo or '').strip().lower()
            if not correo_l or correo_l == yo_l or correo_l in vistos:
                continue
            vistos.add(correo_l)
            destinatarios.append(formataddr((nombre, correo)))

        para = ', '.join(destinatarios)
        asunto = asunto_orig if asunto_orig.lower().startswith('re:') else 'Re: ' + asunto_orig
        msg = MIMEText(texto, 'plain', 'utf-8')
        msg['To'] = para
        msg['Subject'] = asunto
        if msg_id_hdr:
            msg['In-Reply-To'] = msg_id_hdr
            msg['References'] = (refs + ' ' + msg_id_hdr) if refs else msg_id_hdr
        raw = base64.urlsafe_b64encode(msg.as_bytes()).decode()
        enviado = svc.users().messages().send(
            userId='me', body={'raw': raw, 'threadId': thread_id}).execute()
        _audit("gmail_responder_a_todos",
               {"message_id": message_id, "para": para, "asunto": asunto, "id": enviado.get('id')})
        return {"ok": True, "id": enviado.get('id'), "threadId": enviado.get('threadId'), "para": para}
    except Exception as e:
        return _err("gmail_responder_a_todos", e)


@mcp.tool()
def gmail_reenviar(message_id: str, para: str, comentario: str = ''):
    """Reenvia un correo a un nuevo destinatario.

    Recupera el mensaje original (asunto, cabeceras From/To/Date y cuerpo de
    texto plano), compone un mensaje nuevo dirigido a 'para' con asunto 'Fwd:'
    y cuerpo formado por el comentario opcional seguido de una cabecera de
    reenvio y del texto original. No re-adjunta binarios; si el original tenia
    adjuntos se indica en el cuerpo. Devuelve el id del mensaje enviado.
    """
    try:
        from email.mime.text import MIMEText
        svc = _get_gmail()
        m = svc.users().messages().get(userId='me', id=message_id, format='full').execute()
        payload = m.get('payload', {})
        headers = payload.get('headers', [])

        def _h(name):
            for h in headers:
                if h.get('name', '').lower() == name.lower():
                    return h.get('value')
            return ''

        asunto_orig = _h('Subject')

        def _walk(part):
            piezas = [part]
            for sub in part.get('parts', []) or []:
                piezas.extend(_walk(sub))
            return piezas

        texto = ''
        tiene_adjuntos = False
        for p in _walk(payload):
            fn = p.get('filename')
            body = p.get('body', {})
            if fn:
                tiene_adjuntos = True
            elif p.get('mimeType') == 'text/plain' and body.get('data'):
                texto += base64.urlsafe_b64decode(body['data']).decode('utf-8', 'replace')

        cabecera_orig = (
            "De: {de}\n"
            "Fecha: {fecha}\n"
            "Para: {a}\n"
            "Asunto: {asunto}\n"
        ).format(de=_h('From'), fecha=_h('Date'), a=_h('To'), asunto=asunto_orig)

        cuerpo = ''
        if comentario:
            cuerpo += comentario + '\n\n'
        cuerpo += '---------- Mensaje reenviado ----------\n'
        cuerpo += cabecera_orig + '\n' + texto
        if tiene_adjuntos:
            cuerpo += '\n\n[El mensaje original contenia adjuntos que no se han reenviado.]'

        asunto = asunto_orig if asunto_orig.lower().startswith('fwd:') else 'Fwd: ' + asunto_orig
        msg = MIMEText(cuerpo, 'plain', 'utf-8')
        msg['To'] = para
        msg['Subject'] = asunto
        raw = base64.urlsafe_b64encode(msg.as_bytes()).decode()
        enviado = svc.users().messages().send(userId='me', body={'raw': raw}).execute()
        _audit("gmail_reenviar",
               {"message_id": message_id, "para": para, "asunto": asunto, "id": enviado.get('id')})
        return {"ok": True, "id": enviado.get('id'), "threadId": enviado.get('threadId')}
    except Exception as e:
        return _err("gmail_reenviar", e)


@mcp.tool()
def gmail_a_papelera(message_id: str):
    """Mueve un correo a la papelera de Gmail (messages().trash). Devuelve el id y las etiquetas resultantes."""
    try:
        svc = _get_gmail()
        res = svc.users().messages().trash(userId='me', id=message_id).execute()
        _audit("gmail_a_papelera", {"message_id": message_id})
        return {"ok": True, "id": res.get('id'), "labelIds": res.get('labelIds')}
    except Exception as e:
        return _err("gmail_a_papelera", e)


@mcp.tool()
def gmail_restaurar_de_papelera(message_id: str):
    """Restaura un correo desde la papelera de Gmail (messages().untrash). Devuelve el id y las etiquetas resultantes."""
    try:
        svc = _get_gmail()
        res = svc.users().messages().untrash(userId='me', id=message_id).execute()
        _audit("gmail_restaurar_de_papelera", {"message_id": message_id})
        return {"ok": True, "id": res.get('id'), "labelIds": res.get('labelIds')}
    except Exception as e:
        return _err("gmail_restaurar_de_papelera", e)


@mcp.tool()
def gmail_eliminar_definitivo(message_id: str, confirm: bool = False):
    """Elimina un correo de forma DEFINITIVA e irreversible (messages().delete).

    Por seguridad, si confirm es False no borra nada y devuelve un error pidiendo
    volver a llamar con confirm=True. Con confirm=True elimina el mensaje sin
    posibilidad de recuperacion.
    """
    try:
        if not confirm:
            return {"ok": False, "error": "Accion irreversible: vuelve a llamar con confirm=True para eliminar definitivamente el mensaje."}
        svc = _get_gmail()
        svc.users().messages().delete(userId='me', id=message_id).execute()
        _audit("gmail_eliminar_definitivo", {"message_id": message_id})
        return {"ok": True, "id": message_id, "eliminado": True}
    except Exception as e:
        return _err("gmail_eliminar_definitivo", e)


@mcp.tool()
def gmail_destacar(message_id: str, destacar: bool = True):
    """Destaca o quita el destacado de un correo (etiqueta 'STARRED') mediante messages().modify.

    Si destacar es True anade la etiqueta STARRED; si es False la elimina.
    """
    try:
        svc = _get_gmail()
        if destacar:
            body = {'addLabelIds': ['STARRED']}
        else:
            body = {'removeLabelIds': ['STARRED']}
        res = svc.users().messages().modify(userId='me', id=message_id, body=body).execute()
        _audit("gmail_destacar", {"message_id": message_id, "destacar": destacar})
        return {"ok": True, "id": res.get('id'), "labelIds": res.get('labelIds')}
    except Exception as e:
        return _err("gmail_destacar", e)


@mcp.tool()
def gmail_marcar_spam(message_id: str, es_spam: bool = True):
    """Marca o desmarca un correo como spam mediante messages().modify.

    Si es_spam es True anade la etiqueta 'SPAM' y quita 'INBOX'; si es False
    quita 'SPAM' (devolviendo el mensaje fuera de la carpeta de spam).
    """
    try:
        svc = _get_gmail()
        if es_spam:
            body = {'addLabelIds': ['SPAM'], 'removeLabelIds': ['INBOX']}
        else:
            body = {'removeLabelIds': ['SPAM']}
        res = svc.users().messages().modify(userId='me', id=message_id, body=body).execute()
        _audit("gmail_marcar_spam", {"message_id": message_id, "es_spam": es_spam})
        return {"ok": True, "id": res.get('id'), "labelIds": res.get('labelIds')}
    except Exception as e:
        return _err("gmail_marcar_spam", e)


@mcp.tool()
def gmail_quitar_de_carpeta(message_id: str, carpeta: str):
    """Quita un correo de una carpeta/etiqueta identificada por su NOMBRE.

    Busca la etiqueta por nombre en labels().list; si existe, la elimina del
    mensaje con messages().modify (removeLabelIds). Si no existe la etiqueta,
    devuelve un error indicandolo.
    """
    try:
        svc = _get_gmail()
        etiquetas = svc.users().labels().list(userId='me').execute().get('labels', [])
        label_id = None
        for lab in etiquetas:
            if (lab.get('name', '') or '').lower() == carpeta.strip().lower():
                label_id = lab.get('id')
                break
        if not label_id:
            return {"ok": False, "error": "No existe ninguna etiqueta llamada '{}'.".format(carpeta)}
        res = svc.users().messages().modify(
            userId='me', id=message_id, body={'removeLabelIds': [label_id]}).execute()
        _audit("gmail_quitar_de_carpeta",
               {"message_id": message_id, "carpeta": carpeta, "label_id": label_id})
        return {"ok": True, "id": res.get('id'), "labelIds": res.get('labelIds')}
    except Exception as e:
        return _err("gmail_quitar_de_carpeta", e)



# ============================================================
# BLOQUE: Gmail avanzado (PDF, adjuntos, fuera de oficina)
# ============================================================

@mcp.tool()
def gmail_guardar_como_pdf(message_id: str, parent_id: str = 'root', incluir_hilo: bool = False, nombre: Optional[str] = None) -> str:
    """Guarda un correo de Gmail (o el hilo completo) como PDF en Google Drive.

    Obtiene el mensaje indicado o, si `incluir_hilo` es True, todos los mensajes del
    hilo. De cada mensaje extrae las cabeceras From, To, Date y Subject junto con el
    texto plano (text/plain), genera un PDF con reportlab (canvas sobre A4, con saltos
    de linea y control de fin de pagina) y lo sube a la carpeta `parent_id` de Drive.
    El nombre sera `nombre` o el asunto del primer mensaje mas '.pdf'. Devuelve el id
    del fichero creado en Drive.
    """
    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import A4
        import base64
        import textwrap

        svc = _get_gmail()

        def _leer_cabecera(headers, nombre_h):
            for h in headers or []:
                if h.get('name', '').lower() == nombre_h.lower():
                    return h.get('value', '')
            return ''

        def _extraer_texto(payload):
            partes_texto = []

            def _recorrer(part):
                mime = part.get('mimeType', '') or ''
                body = part.get('body', {}) or {}
                data = body.get('data')
                if mime == 'text/plain' and data:
                    try:
                        partes_texto.append(
                            base64.urlsafe_b64decode(data.encode('utf-8')).decode('utf-8', errors='replace')
                        )
                    except Exception:
                        pass
                for sub in part.get('parts', []) or []:
                    _recorrer(sub)

            _recorrer(payload)
            return '\n'.join(partes_texto)

        if incluir_hilo:
            base_msg = svc.users().messages().get(userId='me', id=message_id, format='full').execute()
            thread_id = base_msg.get('threadId')
            hilo = svc.users().threads().get(userId='me', id=thread_id, format='full').execute()
            mensajes = hilo.get('messages', []) or []
        else:
            mensajes = [svc.users().messages().get(userId='me', id=message_id, format='full').execute()]

        primer_payload = (mensajes[0].get('payload') if mensajes else {}) or {}
        asunto = _leer_cabecera(primer_payload.get('headers'), 'Subject') or 'correo'

        buffer = io.BytesIO()
        c = canvas.Canvas(buffer, pagesize=A4)
        ancho, alto = A4
        margen = 50
        y = alto - margen

        def _nueva_pagina():
            nonlocal y
            c.showPage()
            y = alto - margen

        def _escribir_linea(texto, fuente='Helvetica', tam=10):
            nonlocal y
            if y < margen:
                _nueva_pagina()
            c.setFont(fuente, tam)
            c.drawString(margen, y, texto)
            y -= (tam + 4)

        for idx, m in enumerate(mensajes):
            payload = m.get('payload', {}) or {}
            headers = payload.get('headers', [])
            if idx > 0:
                _escribir_linea('-' * 80)
                _escribir_linea('')
            _escribir_linea('De: ' + _leer_cabecera(headers, 'From'), 'Helvetica-Bold', 10)
            _escribir_linea('Para: ' + _leer_cabecera(headers, 'To'), 'Helvetica-Bold', 10)
            _escribir_linea('Fecha: ' + _leer_cabecera(headers, 'Date'), 'Helvetica-Bold', 10)
            _escribir_linea('Asunto: ' + _leer_cabecera(headers, 'Subject'), 'Helvetica-Bold', 10)
            _escribir_linea('')
            cuerpo = _extraer_texto(payload)
            for parrafo in cuerpo.split('\n'):
                if not parrafo:
                    _escribir_linea('')
                    continue
                for linea in textwrap.wrap(parrafo, width=95):
                    _escribir_linea(linea, 'Helvetica', 10)
            _escribir_linea('')

        c.showPage()
        c.save()
        buffer.seek(0)

        nombre_final = nombre or (asunto + '.pdf')
        if not nombre_final.lower().endswith('.pdf'):
            nombre_final = nombre_final + '.pdf'

        drive = _get_service()
        media = MediaIoBaseUpload(buffer, mimetype='application/pdf', resumable=False)
        creado = drive.files().create(
            body={'name': nombre_final, 'parents': [parent_id]},
            media_body=media,
            fields='id',
            supportsAllDrives=True,
        ).execute()
        file_id = creado.get('id')
        _audit('gmail_guardar_como_pdf', 'message=%s hilo=%s -> drive=%s' % (message_id, incluir_hilo, file_id))
        return file_id
    except Exception as e:
        return _err('gmail_guardar_como_pdf', e)


@mcp.tool()
def gmail_leer_adjunto(message_id: str, nombre_adjunto: str, max_chars: int = 20000) -> str:
    """Lee el contenido textual de un adjunto de un correo de Gmail sin guardarlo en Drive.

    Recorre las partes del mensaje buscando el adjunto cuyo filename coincide con (o
    contiene) `nombre_adjunto`, lo descarga con attachments().get y, si es PDF, extrae
    el texto con pypdf; si es texto, lo decodifica. Devuelve el texto (truncado a
    `max_chars`). Si el adjunto es binario no legible, lo indica.
    """
    try:
        import base64

        svc = _get_gmail()
        msg = svc.users().messages().get(userId='me', id=message_id, format='full').execute()
        payload = msg.get('payload', {}) or {}

        encontrado = {}

        def _buscar(part):
            if encontrado:
                return
            filename = part.get('filename', '') or ''
            if filename and (
                nombre_adjunto.lower() == filename.lower() or nombre_adjunto.lower() in filename.lower()
            ):
                encontrado['filename'] = filename
                encontrado['mimeType'] = part.get('mimeType', '') or ''
                encontrado['body'] = part.get('body', {}) or {}
                return
            for sub in part.get('parts', []) or []:
                _buscar(sub)

        _buscar(payload)

        if not encontrado:
            return 'No se encontro ningun adjunto que coincida con: %s' % nombre_adjunto

        body = encontrado['body']
        data = body.get('data')
        if not data and body.get('attachmentId'):
            adj = svc.users().messages().attachments().get(
                userId='me', messageId=message_id, id=body['attachmentId']
            ).execute()
            data = adj.get('data')

        if not data:
            return 'El adjunto "%s" no contiene datos descargables.' % encontrado['filename']

        contenido = base64.urlsafe_b64decode(data.encode('utf-8'))
        filename = encontrado['filename']
        mime = encontrado['mimeType']

        texto = ''
        if filename.lower().endswith('.pdf') or 'pdf' in mime.lower():
            from pypdf import PdfReader
            lector = PdfReader(io.BytesIO(contenido))
            paginas = []
            for pag in lector.pages:
                try:
                    paginas.append(pag.extract_text() or '')
                except Exception:
                    pass
            texto = '\n'.join(paginas)
        else:
            try:
                texto = contenido.decode('utf-8')
            except Exception:
                try:
                    texto = contenido.decode('latin-1')
                except Exception:
                    texto = ''

        if not texto.strip():
            return 'El adjunto "%s" (%s) es binario o no legible como texto.' % (filename, mime)

        if len(texto) > max_chars:
            texto = texto[:max_chars] + '\n[...truncado...]'
        return texto
    except Exception as e:
        return _err('gmail_leer_adjunto', e)


@mcp.tool()
def gmail_fuera_de_oficina(activar: bool, asunto: str = '', mensaje: str = '', desde: Optional[str] = None, hasta: Optional[str] = None) -> str:
    """Activa o desactiva la respuesta automatica de vacaciones (fuera de oficina) en Gmail.

    Si `activar` es True, habilita enableAutoReply con responseSubject (`asunto`),
    responseBodyPlainText (`mensaje`) y, si se aportan fechas AAAA-MM-DD (`desde`/`hasta`),
    startTime/endTime en milisegundos epoch. Si `activar` es False, deshabilita la
    respuesta automatica. Devuelve el estado resultante.
    """
    try:
        import datetime

        svc = _get_gmail()
        body = {'enableAutoReply': bool(activar)}

        if activar:
            body['responseSubject'] = asunto
            body['responseBodyPlainText'] = mensaje

            def _a_ms_inicio(fecha):
                dt = datetime.datetime.strptime(fecha, '%Y-%m-%d').replace(tzinfo=datetime.timezone.utc)
                return int(dt.timestamp() * 1000)

            if desde:
                body['startTime'] = _a_ms_inicio(desde)
            if hasta:
                dt_fin = datetime.datetime.strptime(hasta, '%Y-%m-%d').replace(tzinfo=datetime.timezone.utc)
                dt_fin = dt_fin + datetime.timedelta(days=1) - datetime.timedelta(milliseconds=1)
                body['endTime'] = int(dt_fin.timestamp() * 1000)

        resultado = svc.users().settings().updateVacation(userId='me', body=body).execute()
        _audit('gmail_fuera_de_oficina', 'activar=%s desde=%s hasta=%s' % (activar, desde, hasta))
        estado = 'activada' if resultado.get('enableAutoReply') else 'desactivada'
        return 'Respuesta automatica %s. Asunto: %s' % (estado, resultado.get('responseSubject', ''))
    except Exception as e:
        return _err('gmail_fuera_de_oficina', e)


@mcp.tool()
def gmail_estado_fuera_oficina() -> str:
    """Consulta el estado de la respuesta automatica de vacaciones (fuera de oficina) en Gmail.

    Lee settings().getVacation y devuelve si esta activa, junto con el asunto y el
    mensaje configurados.
    """
    try:
        svc = _get_gmail()
        vac = svc.users().settings().getVacation(userId='me').execute()
        activo = bool(vac.get('enableAutoReply'))
        return 'Fuera de oficina: %s\nAsunto: %s\nMensaje: %s' % (
            'ACTIVO' if activo else 'INACTIVO',
            vac.get('responseSubject', ''),
            vac.get('responseBodyPlainText', ''),
        )
    except Exception as e:
        return _err('gmail_estado_fuera_oficina', e)


# =========================================================================
# BLOQUE: Calendar avanzado (usa _get_cal()). Se APPENDEA a server.py.
# No incluye imports de nivel modulo ni redefine helpers existentes.
# =========================================================================


@mcp.tool()
def calendar_crear_evento_recurrente(summary: str, start: str, end: str, recurrencia: str,
                                     veces: Optional[int] = None, hasta: Optional[str] = None,
                                     descripcion: str = '', attendees: Optional[list] = None,
                                     con_meet: bool = False, calendar_id: str = 'primary'):
    """Crea un evento recurrente en Google Calendar.

    Construye la RRULE a partir de `recurrencia` ('DIARIA'|'SEMANAL'|'MENSUAL'|'ANUAL'),
    anadiendo COUNT=`veces` o UNTIL=`hasta` (formato AAAAMMDD). Los parametros `start` y `end`
    van en RFC3339. Si `con_meet` es True se anade una videollamada de Google Meet.
    Devuelve el id del evento y el enlace (htmlLink).
    """
    import datetime
    import uuid
    try:
        cal = _get_cal()
        mapa = {'DIARIA': 'DAILY', 'SEMANAL': 'WEEKLY', 'MENSUAL': 'MONTHLY', 'ANUAL': 'YEARLY'}
        freq = mapa.get(recurrencia.upper())
        if not freq:
            return _err('calendar_crear_evento_recurrente',
                        ValueError("recurrencia debe ser DIARIA|SEMANAL|MENSUAL|ANUAL"))
        rrule = "RRULE:FREQ=" + freq
        if veces is not None:
            rrule += ";COUNT=" + str(int(veces))
        elif hasta:
            rrule += ";UNTIL=" + str(hasta)
        cuerpo = {
            'summary': summary,
            'description': descripcion,
            'start': {'dateTime': start},
            'end': {'dateTime': end},
            'recurrence': [rrule],
        }
        if attendees:
            cuerpo['attendees'] = [{'email': a} for a in attendees]
        params = {'calendarId': calendar_id, 'body': cuerpo}
        if con_meet:
            cuerpo['conferenceData'] = {
                'createRequest': {
                    'requestId': str(uuid.uuid4()),
                    'conferenceSolutionKey': {'type': 'hangoutsMeet'},
                }
            }
            params['conferenceDataVersion'] = 1
        ev = cal.events().insert(**params).execute()
        _audit('calendar_crear_evento_recurrente',
               "id=%s summary=%s rrule=%s" % (ev.get('id'), summary, rrule))
        return {'id': ev.get('id'), 'link': ev.get('htmlLink')}
    except Exception as e:
        return _err('calendar_crear_evento_recurrente', e)


@mcp.tool()
def calendar_crear_evento_todo_el_dia(summary: str, fecha: str, fecha_fin: Optional[str] = None,
                                      descripcion: str = '', calendar_id: str = 'primary'):
    """Crea un evento de dia completo (all-day) en Google Calendar.

    Usa start={'date': `fecha`} y end={'date': `fecha_fin`}. Si `fecha_fin` no se indica,
    se toma el dia siguiente a `fecha` (las fechas van en formato AAAA-MM-DD).
    Devuelve el id del evento y el enlace (htmlLink).
    """
    import datetime
    import uuid
    try:
        cal = _get_cal()
        if fecha_fin:
            fin = fecha_fin
        else:
            d = datetime.date.fromisoformat(fecha)
            fin = (d + datetime.timedelta(days=1)).isoformat()
        cuerpo = {
            'summary': summary,
            'description': descripcion,
            'start': {'date': fecha},
            'end': {'date': fin},
        }
        ev = cal.events().insert(calendarId=calendar_id, body=cuerpo).execute()
        _audit('calendar_crear_evento_todo_el_dia',
               "id=%s summary=%s %s->%s" % (ev.get('id'), summary, fecha, fin))
        return {'id': ev.get('id'), 'link': ev.get('htmlLink')}
    except Exception as e:
        return _err('calendar_crear_evento_todo_el_dia', e)


@mcp.tool()
def calendar_listar_calendarios():
    """Lista los calendarios disponibles del usuario.

    Recorre calendarList().list() y devuelve, por cada calendario, su id, el resumen
    (summary) y si es el calendario principal (primary).
    """
    import datetime
    import uuid
    try:
        cal = _get_cal()
        salida = []
        page_token = None
        while True:
            resp = cal.calendarList().list(pageToken=page_token).execute()
            for c in resp.get('items', []):
                salida.append({
                    'id': c.get('id'),
                    'resumen': c.get('summary'),
                    'principal': bool(c.get('primary', False)),
                })
            page_token = resp.get('nextPageToken')
            if not page_token:
                break
        return salida
    except Exception as e:
        return _err('calendar_listar_calendarios', e)


@mcp.tool()
def calendar_anadir_meet(event_id: str, calendar_id: str = 'primary'):
    """Anade una videollamada de Google Meet a un evento existente.

    Hace patch del evento con conferenceData (createRequest, conferenceSolutionKey
    hangoutsMeet y requestId uuid) usando conferenceDataVersion=1.
    Devuelve el enlace de Meet asociado al evento.
    """
    import datetime
    import uuid
    try:
        cal = _get_cal()
        cuerpo = {
            'conferenceData': {
                'createRequest': {
                    'requestId': str(uuid.uuid4()),
                    'conferenceSolutionKey': {'type': 'hangoutsMeet'},
                }
            }
        }
        ev = cal.events().patch(calendarId=calendar_id, eventId=event_id,
                                body=cuerpo, conferenceDataVersion=1).execute()
        enlace = ev.get('hangoutLink')
        if not enlace:
            for ep in ev.get('conferenceData', {}).get('entryPoints', []):
                if ep.get('entryPointType') == 'video':
                    enlace = ep.get('uri')
                    break
        _audit('calendar_anadir_meet', "event_id=%s meet=%s" % (event_id, enlace))
        return {'event_id': event_id, 'meet': enlace}
    except Exception as e:
        return _err('calendar_anadir_meet', e)


@mcp.tool()
def calendar_proponer_huecos(fecha: str, duracion_min: int = 30, hora_ini: str = '09:00',
                             hora_fin: str = '18:00', calendar_id: str = 'primary'):
    """Propone huecos libres en la agenda para un dia concreto.

    Consulta freebusy en `fecha` (AAAA-MM-DD) entre `hora_ini` y `hora_fin` (horas locales),
    calcula los tramos LIBRES de al menos `duracion_min` minutos y devuelve hasta 3 huecos
    propuestos, cada uno con inicio y fin en RFC3339.
    """
    import datetime
    import uuid
    try:
        cal = _get_cal()
        tz = datetime.datetime.now().astimezone().tzinfo
        d = datetime.date.fromisoformat(fecha)
        hi_h, hi_m = [int(x) for x in hora_ini.split(':')]
        hf_h, hf_m = [int(x) for x in hora_fin.split(':')]
        ini = datetime.datetime(d.year, d.month, d.day, hi_h, hi_m, tzinfo=tz)
        fin = datetime.datetime(d.year, d.month, d.day, hf_h, hf_m, tzinfo=tz)

        body = {
            'timeMin': ini.isoformat(),
            'timeMax': fin.isoformat(),
            'items': [{'id': calendar_id}],
        }
        fb = cal.freebusy().query(body=body).execute()
        ocupados = fb.get('calendars', {}).get(calendar_id, {}).get('busy', [])

        def _parse(v):
            return datetime.datetime.fromisoformat(v.replace('Z', '+00:00')).astimezone(tz)

        tramos = []
        for b in ocupados:
            bs = max(_parse(b['start']), ini)
            be = min(_parse(b['end']), fin)
            if be > bs:
                tramos.append((bs, be))
        tramos.sort(key=lambda t: t[0])

        # Fusiona solapamientos
        fusion = []
        for bs, be in tramos:
            if fusion and bs <= fusion[-1][1]:
                fusion[-1] = (fusion[-1][0], max(fusion[-1][1], be))
            else:
                fusion.append((bs, be))

        dur = datetime.timedelta(minutes=int(duracion_min))
        huecos = []
        cursor = ini
        for bs, be in fusion:
            if bs - cursor >= dur:
                huecos.append({'inicio': cursor.isoformat(),
                               'fin': (cursor + dur).isoformat()})
            if be > cursor:
                cursor = be
            if len(huecos) >= 3:
                break
        if len(huecos) < 3 and fin - cursor >= dur:
            huecos.append({'inicio': cursor.isoformat(),
                           'fin': (cursor + dur).isoformat()})
        return huecos[:3]
    except Exception as e:
        return _err('calendar_proponer_huecos', e)




# ======================================================================
# BLOQUE: Google Tasks avanzado (helpers _get_tasks / _err / _audit)
# ======================================================================


@mcp.tool()
def tasks_listas() -> dict:
    """Lista todas las listas de tareas de Google Tasks.

    Devuelve el id y el titulo de cada lista de tareas del usuario.
    """
    try:
        service = _get_tasks()
        resultado = service.tasklists().list().execute()
        listas = [
            {"id": item.get("id"), "titulo": item.get("title")}
            for item in resultado.get("items", [])
        ]
        return {"ok": True, "listas": listas}
    except Exception as e:
        return _err("tasks_listas", e)


@mcp.tool()
def tasks_crear_lista(titulo: str) -> dict:
    """Crea una nueva lista de tareas en Google Tasks.

    Args:
        titulo: Titulo de la nueva lista de tareas.

    Devuelve el id de la lista creada.
    """
    try:
        service = _get_tasks()
        creada = service.tasklists().insert(body={"title": titulo}).execute()
        list_id = creada.get("id")
        _audit("tasks_crear_lista", f"titulo={titulo} id={list_id}")
        return {"ok": True, "id": list_id, "titulo": creada.get("title")}
    except Exception as e:
        return _err("tasks_crear_lista", e)


@mcp.tool()
def tasks_editar(
    task_id: str,
    titulo: Optional[str] = None,
    fecha: Optional[str] = None,
    notas: Optional[str] = None,
    tasklist: str = "@default",
) -> dict:
    """Edita una tarea existente de Google Tasks.

    Aplica un patch con los campos no nulos. La fecha en formato
    AAAA-MM-DD se convierte a 'due' RFC3339 (T00:00:00.000Z).

    Args:
        task_id: Identificador de la tarea a editar.
        titulo: Nuevo titulo (opcional).
        fecha: Nueva fecha de vencimiento AAAA-MM-DD (opcional).
        notas: Nuevas notas (opcional).
        tasklist: Lista de tareas donde reside la tarea.

    Devuelve la tarea actualizada.
    """
    try:
        service = _get_tasks()
        body: dict = {}
        if titulo is not None:
            body["title"] = titulo
        if notas is not None:
            body["notes"] = notas
        if fecha is not None:
            body["due"] = fecha + "T00:00:00.000Z"
        actualizada = service.tasks().patch(
            tasklist=tasklist, task=task_id, body=body
        ).execute()
        _audit("tasks_editar", f"task_id={task_id} campos={list(body.keys())}")
        return {"ok": True, "tarea": actualizada}
    except Exception as e:
        return _err("tasks_editar", e)


@mcp.tool()
def tasks_eliminar(task_id: str, tasklist: str = "@default") -> dict:
    """Elimina una tarea de Google Tasks.

    Args:
        task_id: Identificador de la tarea a eliminar.
        tasklist: Lista de tareas donde reside la tarea.

    Devuelve confirmacion de borrado.
    """
    try:
        service = _get_tasks()
        service.tasks().delete(tasklist=tasklist, task=task_id).execute()
        _audit("tasks_eliminar", f"task_id={task_id} tasklist={tasklist}")
        return {"ok": True, "task_id": task_id}
    except Exception as e:
        return _err("tasks_eliminar", e)


@mcp.tool()
def tasks_mover(
    task_id: str,
    tasklist_destino: str,
    tasklist_origen: str = "@default",
) -> dict:
    """Mueve una tarea de una lista de tareas a otra.

    Recupera la tarea de origen, la inserta en la lista de destino
    conservando titulo, notas y fecha de vencimiento, y elimina la
    tarea original.

    Args:
        task_id: Identificador de la tarea a mover.
        tasklist_destino: Lista de tareas de destino.
        tasklist_origen: Lista de tareas de origen.

    Devuelve el nuevo id de la tarea creada en el destino.
    """
    try:
        service = _get_tasks()
        original = service.tasks().get(
            tasklist=tasklist_origen, task=task_id
        ).execute()
        body: dict = {"title": original.get("title", "")}
        if original.get("notes"):
            body["notes"] = original["notes"]
        if original.get("due"):
            body["due"] = original["due"]
        nueva = service.tasks().insert(
            tasklist=tasklist_destino, body=body
        ).execute()
        service.tasks().delete(
            tasklist=tasklist_origen, task=task_id
        ).execute()
        nuevo_id = nueva.get("id")
        _audit(
            "tasks_mover",
            f"task_id={task_id} origen={tasklist_origen} "
            f"destino={tasklist_destino} nuevo_id={nuevo_id}",
        )
        return {"ok": True, "id": nuevo_id}
    except Exception as e:
        return _err("tasks_mover", e)


@mcp.tool()
def tasks_crear_subtarea(
    task_id_padre: str,
    titulo: str,
    notas: str = "",
    tasklist: str = "@default",
) -> dict:
    """Crea una subtarea colgada de una tarea padre en Google Tasks.

    Args:
        task_id_padre: Identificador de la tarea padre.
        titulo: Titulo de la subtarea.
        notas: Notas de la subtarea (opcional).
        tasklist: Lista de tareas donde reside la tarea padre.

    Devuelve el id de la subtarea creada.
    """
    try:
        service = _get_tasks()
        creada = service.tasks().insert(
            tasklist=tasklist,
            parent=task_id_padre,
            body={"title": titulo, "notes": notas},
        ).execute()
        sub_id = creada.get("id")
        _audit(
            "tasks_crear_subtarea",
            f"padre={task_id_padre} id={sub_id} titulo={titulo}",
        )
        return {"ok": True, "id": sub_id}
    except Exception as e:
        return _err("tasks_crear_subtarea", e)


# =============================================================================
# BLOQUE: Contactos ESCRITURA (People API v1) — requiere scope de escritura
# =============================================================================


@mcp.tool()
def contactos_crear(
    nombre: str,
    email: Optional[str] = None,
    telefono: Optional[str] = None,
    empresa: Optional[str] = None,
) -> dict:
    """Crea un contacto nuevo en Google Contacts (People API).

    Parámetros:
        nombre: Nombre a mostrar del contacto (obligatorio).
        email: Dirección de correo electrónico (opcional).
        telefono: Número de teléfono (opcional).
        empresa: Nombre de la organización o empresa (opcional).

    Devuelve un diccionario con el resourceName y el nombre del contacto creado.
    """
    try:
        people = _get_people()
        body: dict = {
            "names": [{"givenName": nombre, "displayName": nombre}],
        }
        if email:
            body["emailAddresses"] = [{"value": email}]
        if telefono:
            body["phoneNumbers"] = [{"value": telefono}]
        if empresa:
            body["organizations"] = [{"name": empresa}]

        resultado = people.people().createContact(body=body).execute()
        resource_name = resultado.get("resourceName")
        _audit("contactos_crear", f"nombre={nombre} resourceName={resource_name}")
        return {"ok": True, "resourceName": resource_name, "nombre": nombre}
    except Exception as e:
        return _err("contactos_crear", e)


@mcp.tool()
def contactos_editar(
    resource_name: str,
    nombre: Optional[str] = None,
    email: Optional[str] = None,
    telefono: Optional[str] = None,
) -> dict:
    """Edita un contacto existente en Google Contacts (People API).

    Obtiene primero el etag del contacto y actualiza únicamente los campos
    indicados (nombre, email y/o teléfono).

    Parámetros:
        resource_name: Identificador del contacto (p. ej. 'people/c123...').
        nombre: Nuevo nombre a mostrar (opcional).
        email: Nueva dirección de correo (opcional).
        telefono: Nuevo número de teléfono (opcional).

    Devuelve un diccionario con ok=True si la actualización se realiza.
    """
    try:
        people = _get_people()
        actual = (
            people.people()
            .get(
                resourceName=resource_name,
                personFields="names,emailAddresses,phoneNumbers,metadata",
            )
            .execute()
        )
        etag = actual.get("etag")

        campos = []
        body: dict = {"etag": etag}
        if nombre is not None:
            body["names"] = [{"givenName": nombre, "displayName": nombre}]
            campos.append("names")
        if email is not None:
            body["emailAddresses"] = [{"value": email}]
            campos.append("emailAddresses")
        if telefono is not None:
            body["phoneNumbers"] = [{"value": telefono}]
            campos.append("phoneNumbers")

        if not campos:
            return _err(
                "contactos_editar",
                ValueError("No se indicó ningún campo a modificar."),
            )

        people.people().updateContact(
            resourceName=resource_name,
            updatePersonFields=",".join(campos),
            body=body,
        ).execute()
        _audit(
            "contactos_editar",
            f"resourceName={resource_name} campos={','.join(campos)}",
        )
        return {"ok": True, "resourceName": resource_name}
    except Exception as e:
        return _err("contactos_editar", e)


@mcp.tool()
def contactos_eliminar(resource_name: str, confirm: bool = False) -> dict:
    """Elimina un contacto de Google Contacts (People API).

    Por seguridad requiere confirmación explícita.

    Parámetros:
        resource_name: Identificador del contacto (p. ej. 'people/c123...').
        confirm: Debe ser True para ejecutar el borrado; si es False se aborta.

    Devuelve un diccionario con ok=True si el contacto se elimina.
    """
    try:
        if not confirm:
            return _err(
                "contactos_eliminar",
                ValueError(
                    "Operación no confirmada: vuelve a llamar con confirm=True "
                    "para eliminar el contacto."
                ),
            )
        people = _get_people()
        people.people().deleteContact(resourceName=resource_name).execute()
        _audit("contactos_eliminar", f"resourceName={resource_name}")
        return {"ok": True, "resourceName": resource_name}
    except Exception as e:
        return _err("contactos_eliminar", e)



# --------------------------------------------------------------------------- #
# COMANDOS COMBINADOS DEL DIA (bandeja y agenda)
# --------------------------------------------------------------------------- #


@mcp.tool()
def bandeja_del_dia(max_correos: int = 15) -> dict:
    """Reune la bandeja del dia: correos de Gmail no leidos y recientes
    (ultimos 2 dias) y tareas de Google Tasks vencidas o que vencen hoy.

    Args:
        max_correos: Numero maximo de correos no leidos a devolver.

    Devuelve un dict con 'correos_no_leidos' (lista con remitente, asunto y
    fecha de cada correo) y 'tareas_pendientes' (tareas con vencimiento <= hoy
    que no esten completadas).
    """
    try:
        import datetime as _dt
        gmail = _get_gmail()
        tasks = _get_tasks()

        # (a) Correos no leidos recientes (ultimos 2 dias)
        correos = []
        res = gmail.users().messages().list(
            userId="me", q="is:unread newer_than:2d",
            maxResults=max(1, min(max_correos, 100)),
        ).execute()
        for m in res.get("messages", []):
            full = gmail.users().messages().get(
                userId="me", id=m["id"], format="metadata",
                metadataHeaders=["From", "Subject", "Date"],
            ).execute()
            h = {x["name"]: x["value"]
                 for x in full.get("payload", {}).get("headers", [])}
            correos.append({
                "id": m["id"],
                "de": h.get("From"),
                "asunto": h.get("Subject"),
                "fecha": h.get("Date"),
                "resumen": full.get("snippet", "")[:160],
            })

        # (b) Tareas de @default vencidas o que vencen hoy (no completadas)
        hoy = _dt.date.today().isoformat()
        pendientes = []
        tres = tasks.tasks().list(
            tasklist="@default", maxResults=100, showCompleted=False,
        ).execute()
        for t in tres.get("items", []):
            if t.get("status") == "completed":
                continue
            due = t.get("due")
            if not due:
                continue
            if due[:10] <= hoy:
                pendientes.append({
                    "id": t.get("id"),
                    "titulo": t.get("title"),
                    "estado": t.get("status"),
                    "vencimiento": due,
                })

        return {
            "ok": True,
            "fecha": hoy,
            "correos_no_leidos": correos,
            "tareas_pendientes": pendientes,
        }
    except Exception as e:
        return _err("bandeja_del_dia", e)


@mcp.tool()
def agenda_del_dia(fecha: Optional[str] = None) -> dict:
    """Agenda del dia: eventos de Google Calendar de la fecha indicada con sus
    documentos relacionados en Drive.

    Para cada evento devuelve hora, titulo, asistentes y enlace de Google Meet.
    Ademas busca en Drive archivos cuyo nombre contenga el titulo del evento o
    el nombre del asistente principal y los adjunta como
    'documentos_relacionados'.

    Args:
        fecha: Dia a consultar en formato 'AAAA-MM-DD'. Si no se indica, hoy.

    Devuelve la lista de eventos del dia con sus documentos.
    """
    try:
        import datetime as _dt
        if not fecha:
            fecha = _dt.date.today().isoformat()
        cal = _get_cal()
        drv = _get_service()

        time_min = fecha + "T00:00:00Z"
        time_max = fecha + "T23:59:59Z"
        res = cal.events().list(
            calendarId="primary", timeMin=time_min, timeMax=time_max,
            singleEvents=True, orderBy="startTime", maxResults=100,
        ).execute()

        eventos = []
        for e in res.get("items", []):
            st = e.get("start", {})
            en = e.get("end", {})
            titulo = e.get("summary") or ""
            asistentes = [a.get("email") for a in e.get("attendees", [])
                          if a.get("email")]

            # Terminos de busqueda: titulo y nombre del asistente principal
            terminos = []
            if titulo.strip():
                terminos.append(titulo.strip())
            if asistentes:
                principal = asistentes[0].split("@")[0].replace(".", " ").strip()
                if principal:
                    terminos.append(principal)

            docs = []
            if terminos:
                clausulas = []
                for term in terminos:
                    seguro = term.replace("\\", " ").replace("'", " ").strip()
                    if seguro:
                        clausulas.append("name contains '%s'" % seguro)
                if clausulas:
                    q = "(" + " or ".join(clausulas) + ") and trashed = false"
                    try:
                        dres = drv.files().list(
                            q=q, pageSize=5,
                            fields="files(%s)" % FILE_FIELDS,
                            supportsAllDrives=True,
                            includeItemsFromAllDrives=True,
                        ).execute()
                        docs = [_file(f) for f in dres.get("files", [])]
                    except Exception:
                        docs = []

            eventos.append({
                "id": e.get("id"),
                "titulo": titulo,
                "hora_inicio": st.get("dateTime") or st.get("date"),
                "hora_fin": en.get("dateTime") or en.get("date"),
                "asistentes": asistentes,
                "meet": e.get("hangoutLink"),
                "link": e.get("htmlLink"),
                "documentos_relacionados": docs,
            })

        return {
            "ok": True,
            "fecha": fecha,
            "count": len(eventos),
            "eventos": eventos,
        }
    except Exception as e:
        return _err("agenda_del_dia", e)


# =========================================================================
# BLOQUE DLP (Prevencion de Fuga de Datos) - solo informa, no modifica nada
# =========================================================================


@mcp.tool()
def dlp_escanear_compartidos(max_resultados: int = 300) -> dict:
    """Escanea Google Drive en busca de ficheros compartidos por enlace publico
    ('cualquiera con el enlace' o 'cualquiera puede encontrarlo') cuyo nombre
    contenga palabras sensibles (token, clave, password, DNI, nomina, IBAN,
    factura, contrato, historia clinica, etc.) y los marca como ALERTA con la
    palabra coincidente. Solo informa: NO cambia permisos ni ficheros.

    Args:
        max_resultados: numero maximo de ficheros a devolver como alerta (por
            defecto 300). Se lista una sola pagina para no agotar el tiempo.

    Returns:
        dict con el total revisado, el total de alertas y la lista de alertas.
    """
    action = "dlp_escanear_compartidos"
    try:
        import re  # noqa: F401
        service = _get_service()
        palabras = [
            "token", "clave", "password", "contrase", "dni", "nie",
            "pasaporte", "nomina", "nómina", "vida laboral", "iban",
            "factura", "contrato", "historia clinica",
        ]
        q = ("(visibility='anyoneWithLink' or visibility='anyoneCanFind') "
             "and trashed=false")
        resp = service.files().list(
            q=q,
            pageSize=1000,
            fields="nextPageToken, files(%s)" % FILE_FIELDS,
            supportsAllDrives=True,
            includeItemsFromAllDrives=True,
        ).execute()
        ficheros = resp.get("files", [])
        alertas = []
        for meta in ficheros:
            nombre_low = (meta.get("name") or "").lower()
            coincidencias = [p for p in palabras if p in nombre_low]
            if coincidencias:
                item = _file(meta)
                item["alerta"] = True
                item["palabras_coincidentes"] = coincidencias
                alertas.append(item)
                if len(alertas) >= max_resultados:
                    break
        _audit(action, "revisados=%d alertas=%d" % (len(ficheros), len(alertas)))
        return {
            "ok": True,
            "total_revisados": len(ficheros),
            "total_alertas": len(alertas),
            "alertas": alertas,
        }
    except Exception as e:
        return _err(action, e)


@mcp.tool()
def dlp_revisar_texto(texto: str) -> dict:
    """Analiza un texto (correo, documento o mensaje) y detecta datos personales
    sensibles mediante expresiones regulares: DNI (8 digitos + letra), NIE
    (X/Y/Z + 7 digitos + letra), IBAN espanol (ES + 22 digitos), telefono
    espanol (9 digitos), email y numero de la Seguridad Social (12 digitos).
    Devuelve que tipos aparecen y cuantos de cada uno, ENMASCARANDO los valores
    (solo se muestran los ultimos 3 caracteres). Util para revisar algo antes de
    enviarlo. Solo informa: NO modifica nada.

    Args:
        texto: texto a inspeccionar antes de enviarlo o publicarlo.

    Returns:
        dict con los tipos de dato sensible detectados, su recuento y muestras
        enmascaradas.
    """
    action = "dlp_revisar_texto"
    try:
        import re
        texto = texto or ""
        patrones = {
            "email": r"\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}\b",
            "IBAN_ES": r"\bES\d{22}\b",
            "NIE": r"\b[XYZxyz]\d{7}[A-Za-z]\b",
            "DNI": r"\b\d{8}[A-Za-z]\b",
            "num_seguridad_social": r"(?<!\d)\d{12}(?!\d)",
            "telefono_ES": r"(?<!\d)\d{9}(?!\d)",
        }
        orden = [
            "email", "IBAN_ES", "NIE", "DNI",
            "num_seguridad_social", "telefono_ES",
        ]
        trabajo = texto
        detectados = {}
        for tipo in orden:
            encontrados = []

            def _rep(m, _acc=encontrados):
                v = m.group(0)
                _acc.append(v)
                return " " * len(v)

            trabajo = re.sub(patrones[tipo], _rep, trabajo)
            if encontrados:
                muestras = []
                for v in encontrados:
                    vs = v.strip()
                    if len(vs) > 3:
                        muestras.append("*" * (len(vs) - 3) + vs[-3:])
                    else:
                        muestras.append("*" * len(vs))
                detectados[tipo] = {
                    "cantidad": len(encontrados),
                    "muestras": muestras,
                }
        _audit(action, "tipos=%s" % ",".join(detectados.keys()))
        return {
            "ok": True,
            "hay_datos_sensibles": bool(detectados),
            "tipos_detectados": list(detectados.keys()),
            "detalle": detectados,
        }
    except Exception as e:
        return _err(action, e)




# =====================================================================
# BLOQUE: Retencion documental (RGPD)
# =====================================================================


@mcp.tool()
def retencion_listar_antiguos(folder_id: str, anios: int = 5,
                              page_size: int = 200) -> dict:
    """Lista los ficheros dentro de `folder_id` cuya fecha de modificacion sea
    anterior a hace `anios` anios (candidatos a retencion/supresion RGPD).
    Excluye los que esten bajo retencion legal (appProperties legal_hold='true').
    Devuelve nombre, fecha e id de cada fichero, mas el total."""
    try:
        from datetime import datetime, timezone, timedelta
        service = _get_service()
        limite = datetime.now(timezone.utc) - timedelta(days=365 * anios)
        q = "'%s' in parents and trashed=false" % folder_id
        fields = ("nextPageToken, files(id, name, mimeType, modifiedTime, "
                  "appProperties)")
        antiguos = []
        page_token = None
        while True:
            res = service.files().list(
                q=q,
                fields=fields,
                pageSize=page_size,
                pageToken=page_token,
                supportsAllDrives=True,
                includeItemsFromAllDrives=True,
            ).execute()
            for f in res.get("files", []):
                props = f.get("appProperties") or {}
                if props.get("legal_hold") == "true":
                    continue
                mtime = f.get("modifiedTime")
                if not mtime:
                    continue
                fecha = datetime.fromisoformat(mtime.replace("Z", "+00:00"))
                if fecha < limite:
                    antiguos.append({
                        "id": f.get("id"),
                        "nombre": f.get("name"),
                        "fecha": mtime,
                    })
            page_token = res.get("nextPageToken")
            if not page_token:
                break
        return {"ok": True, "folder_id": folder_id, "anios": anios,
                "limite": limite.isoformat(), "total": len(antiguos),
                "ficheros": antiguos}
    except Exception as e:
        return _err("retencion_listar_antiguos", e)


@mcp.tool()
def retencion_archivar(file_ids: list, carpeta_archivo_id: str) -> dict:
    """Mueve cada fichero de `file_ids` a la carpeta `carpeta_archivo_id`
    (patron drive_move). Salta los ficheros que esten bajo retencion legal
    (appProperties legal_hold='true'). Devuelve el resultado por fichero."""
    try:
        service = _get_service()
        resultados = []
        for file_id in file_ids:
            try:
                current = service.files().get(
                    fileId=file_id,
                    fields="parents, name, appProperties",
                    supportsAllDrives=True,
                ).execute()
                props = current.get("appProperties") or {}
                if props.get("legal_hold") == "true":
                    resultados.append({
                        "file_id": file_id,
                        "nombre": current.get("name"),
                        "ok": False,
                        "saltado": True,
                        "motivo": "legal_hold activo",
                    })
                    continue
                prev_parents = ",".join(current.get("parents", []))
                meta = service.files().update(
                    fileId=file_id,
                    addParents=carpeta_archivo_id,
                    removeParents=prev_parents,
                    fields=FILE_FIELDS,
                    supportsAllDrives=True,
                ).execute()
                _audit("retencion_archivar", {
                    "file_id": file_id,
                    "nombre": meta.get("name"),
                    "desde": prev_parents,
                    "hacia": carpeta_archivo_id,
                })
                resultados.append({
                    "file_id": file_id,
                    "ok": True,
                    "saltado": False,
                    "file": _file(meta),
                })
            except Exception as inner:
                resultados.append({
                    "file_id": file_id,
                    "ok": False,
                    "saltado": False,
                    "error": str(inner),
                })
        movidos = sum(1 for r in resultados if r.get("ok"))
        saltados = sum(1 for r in resultados if r.get("saltado"))
        return {"ok": True, "carpeta_archivo_id": carpeta_archivo_id,
                "total": len(file_ids), "movidos": movidos,
                "saltados": saltados, "resultados": resultados}
    except Exception as e:
        return _err("retencion_archivar", e)


@mcp.tool()
def retencion_legal_hold(file_id: str, activar: bool = True) -> dict:
    """Marca o desmarca un fichero como bajo retencion legal (legal hold),
    fijando appProperties {'legal_hold': 'true'|'false'}. Los ficheros con
    retencion legal quedan excluidos de la supresion y del archivado RGPD.
    Devuelve el estado resultante."""
    try:
        from datetime import datetime, timezone
        service = _get_service()
        valor = "true" if activar else "false"
        body = {"appProperties": {"legal_hold": valor}}
        meta = service.files().update(
            fileId=file_id,
            body=body,
            fields="id, appProperties",
            supportsAllDrives=True,
        ).execute()
        _audit("retencion_legal_hold", {
            "file_id": file_id,
            "legal_hold": valor,
            "cuando": datetime.now(timezone.utc).isoformat(),
        })
        props = meta.get("appProperties") or {}
        return {"ok": True, "file_id": meta.get("id"),
                "legal_hold": props.get("legal_hold"),
                "activo": props.get("legal_hold") == "true"}
    except Exception as e:
        return _err("retencion_legal_hold", e)


# --------------------------------------------------------------------------- #
# COMPARTICION CON CADUCIDAD (portal ligero)
# =========================================================================== #
# GOOGLE FORMS (requiere scopes forms.body / forms.responses.readonly)
# =========================================================================== #

_forms_service = None

def _get_forms():
    global _forms_service
    if _forms_service is None:
        _forms_service = build("forms", "v1", credentials=_build_creds(),
                               cache_discovery=False,
                               discoveryServiceUrl="https://forms.googleapis.com/$discovery/rest?version=v1")
    return _forms_service


@mcp.tool()
def forms_crear(titulo: str, descripcion: str = "", preguntas: Optional[list] = None) -> dict:
    """Crea un formulario de Google Forms. `preguntas` es una lista de dicts, p. ej.:
    [{"titulo":"Nombre y apellidos","tipo":"texto","requerida":true},
     {"titulo":"Motivo de la consulta","tipo":"parrafo"},
     {"titulo":"Materia","tipo":"opcion","opciones":["Despido","Nomina","Otro"]},
     {"titulo":"Documentos que aporta","tipo":"casillas","opciones":["Contrato","Nominas"]}].
    Tipos: 'texto' (respuesta corta), 'parrafo', 'opcion' (una), 'casillas' (varias).
    Devuelve el id, el enlace para responder y el de edicion."""
    try:
        svc = _get_forms()
        form = svc.forms().create(body={"info": {"title": titulo}}).execute()
        fid = form["formId"]
        requests = []
        if descripcion:
            requests.append({"updateFormInfo": {
                "info": {"description": descripcion}, "updateMask": "description"}})
        for i, p in enumerate(preguntas or []):
            tipo = (p.get("tipo") or "texto").lower()
            req = bool(p.get("requerida"))
            if tipo in ("opcion", "casillas"):
                qtype = "RADIO" if tipo == "opcion" else "CHECKBOX"
                question = {"required": req, "choiceQuestion": {
                    "type": qtype,
                    "options": [{"value": str(o)} for o in (p.get("opciones") or [])]}}
            else:
                question = {"required": req, "textQuestion": {
                    "paragraph": tipo == "parrafo"}}
            requests.append({"createItem": {
                "item": {"title": p.get("titulo", "Pregunta"),
                         "questionItem": {"question": question}},
                "location": {"index": i}}})
        if requests:
            svc.forms().batchUpdate(formId=fid, body={"requests": requests}).execute()
        info = svc.forms().get(formId=fid).execute()
        _audit("forms_crear", {"id": fid, "titulo": titulo})
        return {"ok": True, "id": fid,
                "responder_url": info.get("responderUri"),
                "editar_url": "https://docs.google.com/forms/d/" + fid + "/edit"}
    except Exception as e:
        return _err("forms_crear", e)


@mcp.tool()
def forms_leer(form_id: str) -> dict:
    """Devuelve la estructura de un formulario: titulo, descripcion, enlace para
    responder y la lista de preguntas."""
    try:
        svc = _get_forms()
        f = svc.forms().get(formId=form_id).execute()
        preguntas = []
        for it in f.get("items", []):
            q = it.get("questionItem", {}).get("question", {})
            tipo = "parrafo" if q.get("textQuestion", {}).get("paragraph") else (
                   "opcion/casillas" if q.get("choiceQuestion") else "texto")
            preguntas.append({"titulo": it.get("title"), "tipo": tipo})
        return {"id": form_id, "titulo": f.get("info", {}).get("title"),
                "descripcion": f.get("info", {}).get("description"),
                "responder_url": f.get("responderUri"), "preguntas": preguntas}
    except Exception as e:
        return _err("forms_leer", e)


@mcp.tool()
def forms_respuestas(form_id: str, max_resultados: int = 100) -> dict:
    """Lista las respuestas recibidas en un formulario. Devuelve, por cada respuesta,
    la fecha y las contestaciones (pregunta -> valor)."""
    try:
        svc = _get_forms()
        f = svc.forms().get(formId=form_id).execute()
        # mapa itemId/questionId -> titulo de la pregunta
        qmap = {}
        for it in f.get("items", []):
            qid = it.get("questionItem", {}).get("question", {}).get("questionId")
            if qid:
                qmap[qid] = it.get("title")
        res = svc.forms().responses().list(formId=form_id).execute()
        out = []
        for r in res.get("responses", [])[:max_resultados]:
            answers = {}
            for qid, a in (r.get("answers") or {}).items():
                vals = [v.get("value") for v in a.get("textAnswers", {}).get("answers", [])]
                answers[qmap.get(qid, qid)] = ", ".join(vals)
            out.append({"fecha": r.get("lastSubmittedTime"), "respuestas": answers})
        return {"total": len(out), "respuestas": out}
    except Exception as e:
        return _err("forms_respuestas", e)


# --------------------------------------------------------------------------- #

@mcp.tool()
def drive_compartir_temporal(file_id: str, email: str, dias: int = 7,
                             rol: str = "reader") -> dict:
    """Comparte un fichero con un usuario durante un numero limitado de dias.

    Crea un permiso de tipo 'user' sobre `file_id` para el `email` indicado con
    el rol `rol` y una fecha de caducidad (`expirationTime`) igual a AHORA mas
    `dias` dias, en formato RFC3339 con sufijo 'Z'.

    IMPORTANTE: Google Drive solo admite expiracion para los roles
    'reader', 'commenter' y 'writer' (nunca 'owner'). Si se pasa otro rol se
    devuelve un error accionable sin llamar a la API.

    Devuelve el id del permiso creado, la fecha de caducidad efectiva y el
    enlace (webViewLink) del fichero. El acceso puede revocarse antes de tiempo
    con la herramienta existente `drive_remove_permission`.
    """
    try:
        from datetime import datetime, timedelta, timezone
        roles_validos = ("reader", "commenter", "writer")
        if rol not in roles_validos:
            return {
                "ok": False,
                "action": "drive_compartir_temporal",
                "error": "El rol '%s' no admite caducidad." % rol,
                "sugerencia": "La expiracion solo es valida para los roles: "
                              "reader, commenter o writer (no owner).",
            }
        if dias <= 0:
            return {
                "ok": False,
                "action": "drive_compartir_temporal",
                "error": "El numero de dias debe ser mayor que 0.",
                "sugerencia": "Indica cuantos dias durara el acceso (p.ej. 7).",
            }
        service = _get_service()
        caduca_dt = datetime.now(timezone.utc) + timedelta(days=dias)
        expiration = caduca_dt.strftime("%Y-%m-%dT%H:%M:%SZ")
        body = {
            "type": "user",
            "role": rol,
            "emailAddress": email,
            "expirationTime": expiration,
        }
        perm = service.permissions().create(
            fileId=file_id,
            body=body,
            fields="id,expirationTime",
            supportsAllDrives=True,
            sendNotificationEmail=True,
        ).execute()
        meta = service.files().get(
            fileId=file_id, fields="webViewLink", supportsAllDrives=True
        ).execute()
        _audit("drive_compartir_temporal", {
            "file_id": file_id,
            "email": email,
            "rol": rol,
            "dias": dias,
            "permission_id": perm.get("id"),
            "expiration_time": perm.get("expirationTime", expiration),
        })
        return {
            "ok": True,
            "permission_id": perm.get("id"),
            "caduca": perm.get("expirationTime", expiration),
            "enlace": meta.get("webViewLink"),
            "resumen": "Acceso '%s' a %s hasta %s."
                       % (rol, email, perm.get("expirationTime", expiration)),
        }
    except Exception as e:
        return _err("drive_compartir_temporal", e)


@mcp.tool()
def drive_permisos_con_caducidad(file_id: str) -> dict:
    """Lista los permisos de un fichero que tengan fecha de caducidad.

    Recupera todos los permisos de `file_id` y filtra los que incluyen
    `expirationTime`, es decir, los accesos temporales. Devuelve quien tiene
    acceso temporal, con que rol y hasta cuando, para revisar o revocar (esto
    ultimo con la herramienta existente `drive_remove_permission`).
    """
    try:
        service = _get_service()
        resp = service.permissions().list(
            fileId=file_id,
            fields="permissions(id,emailAddress,role,expirationTime)",
            supportsAllDrives=True,
        ).execute()
        permisos = resp.get("permissions", []) or []
        temporales = [
            {
                "permission_id": p.get("id"),
                "email": p.get("emailAddress"),
                "rol": p.get("role"),
                "caduca": p.get("expirationTime"),
            }
            for p in permisos
            if p.get("expirationTime")
        ]
        return {
            "ok": True,
            "file_id": file_id,
            "total": len(temporales),
            "accesos_temporales": temporales,
        }
    except Exception as e:
        return _err("drive_permisos_con_caducidad", e)


# =========================================================================
# BLOQUE: Control horario y facturacion (almacen en un Google Sheet)
# =========================================================================


def _hoja_horas(parent_id: str = 'root') -> str:
    """Localiza (o crea) el Google Sheet 'Registro de horas - Aurea'.

    Busca en Drive una hoja de calculo con ese nombre. Si no existe, la
    crea en ``parent_id`` y escribe la fila de cabecera. Devuelve el
    spreadsheetId.
    """
    service = _get_service()
    q = ("name='Registro de horas - Aurea' and "
         "mimeType='application/vnd.google-apps.spreadsheet' and "
         "trashed=false")
    encontrados = service.files().list(
        q=q,
        fields='files(id, name)',
        pageSize=1,
    ).execute().get('files', [])
    if encontrados:
        return encontrados[0]['id']
    metadata = {
        'name': 'Registro de horas - Aurea',
        'mimeType': 'application/vnd.google-apps.spreadsheet',
    }
    if parent_id and parent_id != 'root':
        metadata['parents'] = [parent_id]
    creado = service.files().create(body=metadata, fields='id').execute()
    sid = creado['id']
    _get_sheets().spreadsheets().values().update(
        spreadsheetId=sid,
        range='A1',
        valueInputOption='USER_ENTERED',
        body={'values': [['Fecha', 'Cliente', 'Asunto', 'Minutos', 'Concepto']]},
    ).execute()
    return sid


@mcp.tool()
def horario_registrar(cliente: str, minutos: int, concepto: str = '',
                      asunto: str = '', fecha: Optional[str] = None,
                      parent_id: str = 'root') -> dict:
    """Registra el tiempo dedicado a un cliente en la hoja de horas.

    Args:
        cliente: Nombre del cliente.
        minutos: Minutos dedicados.
        concepto: Descripcion del concepto o tarea facturable.
        asunto: Asunto o expediente relacionado.
        fecha: Fecha en formato 'AAAA-MM-DD'; si se omite se usa hoy.
        parent_id: Carpeta de Drive donde crear la hoja si no existe.

    Devuelve un dict con 'ok' y el total de filas registradas.
    """
    try:
        import datetime
        sid = _hoja_horas(parent_id)
        f = fecha or datetime.date.today().isoformat()
        fila = [f, cliente, asunto, minutos, concepto]
        resp = _get_sheets().spreadsheets().values().append(
            spreadsheetId=sid,
            range='A:E',
            valueInputOption='USER_ENTERED',
            body={'values': [fila]},
        ).execute()
        _audit('horario_registrar',
               "cliente=%s minutos=%s fecha=%s" % (cliente, minutos, f))
        columna_a = _get_sheets().spreadsheets().values().get(
            spreadsheetId=sid,
            range='A:A',
        ).execute().get('values', [])
        total_filas = max(0, len(columna_a) - 1)
        return {
            'ok': True,
            'spreadsheetId': sid,
            'fila': fila,
            'rango_actualizado': resp.get('updates', {}).get('updatedRange', ''),
            'total_filas': total_filas,
        }
    except Exception as e:
        return _err('horario_registrar', e)


@mcp.tool()
def horario_informe(cliente: Optional[str] = None,
                    mes: Optional[str] = None) -> dict:
    """Genera un informe de dedicacion a partir de la hoja de horas.

    Args:
        cliente: Si se indica, filtra por este cliente.
        mes: Si se indica (formato 'AAAA-MM'), filtra por ese mes comparando
            el inicio de la fecha de cada registro.

    Devuelve el total de minutos y de horas y el desglose por cliente.
    """
    try:
        sid = _hoja_horas()
        filas = _get_sheets().spreadsheets().values().get(
            spreadsheetId=sid,
            range='A2:E',
        ).execute().get('values', [])
        total_min = 0
        registros = 0
        por_cliente: dict = {}
        for fila in filas:
            f = fila[0] if len(fila) > 0 else ''
            c = fila[1] if len(fila) > 1 else ''
            m = fila[3] if len(fila) > 3 else 0
            if cliente and c != cliente:
                continue
            if mes and not str(f).startswith(mes):
                continue
            try:
                texto = str(m).strip().replace(',', '.')
                mins = int(float(texto)) if texto else 0
            except (ValueError, TypeError):
                mins = 0
            total_min += mins
            registros += 1
            por_cliente[c] = por_cliente.get(c, 0) + mins
        return {
            'ok': True,
            'filtro_cliente': cliente,
            'filtro_mes': mes,
            'registros': registros,
            'total_minutos': total_min,
            'total_horas': round(total_min / 60.0, 2),
            'desglose_por_cliente': por_cliente,
        }
    except Exception as e:
        return _err('horario_informe', e)


@mcp.tool()
def horario_estimar_cliente(cliente: str, dias: int = 30) -> dict:
    """Estima orientativamente la dedicacion a un cliente.

    Combina el numero de correos de Gmail con ese cliente y la duracion de
    las reuniones de Calendar que lo mencionen en el periodo indicado.

    Args:
        cliente: Nombre o termino de busqueda del cliente.
        dias: Numero de dias hacia atras a considerar.

    Devuelve el desglose y una estimacion (10 min por correo + duracion de
    reuniones).
    """
    try:
        import datetime
        gmail = _get_gmail()
        q = "%s newer_than:%sd" % (cliente, dias)
        num_correos = 0
        page_token = None
        while True:
            r = gmail.users().messages().list(
                userId='me', q=q, pageToken=page_token, maxResults=500,
            ).execute()
            num_correos += len(r.get('messages', []))
            page_token = r.get('nextPageToken')
            if not page_token:
                break
        cal = _get_cal()
        ahora = datetime.datetime.utcnow()
        time_min = (ahora - datetime.timedelta(days=dias)).isoformat() + 'Z'
        time_max = ahora.isoformat() + 'Z'
        eventos = cal.events().list(
            calendarId='primary',
            q=cliente,
            timeMin=time_min,
            timeMax=time_max,
            singleEvents=True,
            orderBy='startTime',
        ).execute().get('items', [])
        num_reuniones = 0
        min_reuniones = 0
        for ev in eventos:
            ini = ev.get('start', {}).get('dateTime')
            fin = ev.get('end', {}).get('dateTime')
            if not ini or not fin:
                continue
            try:
                di = datetime.datetime.fromisoformat(ini.replace('Z', '+00:00'))
                df = datetime.datetime.fromisoformat(fin.replace('Z', '+00:00'))
                dur = int((df - di).total_seconds() // 60)
            except (ValueError, TypeError):
                continue
            if dur > 0:
                min_reuniones += dur
                num_reuniones += 1
        min_correos = num_correos * 10
        total_min = min_correos + min_reuniones
        return {
            'ok': True,
            'cliente': cliente,
            'dias': dias,
            'num_correos': num_correos,
            'minutos_correos_estimados': min_correos,
            'num_reuniones': num_reuniones,
            'minutos_reuniones': min_reuniones,
            'estimacion_minutos': total_min,
            'estimacion_horas': round(total_min / 60.0, 2),
            'nota': 'Estimacion orientativa: 10 min por correo + duracion de reuniones.',
        }
    except Exception as e:
        return _err('horario_estimar_cliente', e)




@mcp.tool()
def buscar_en_todo(texto: str, max_por_fuente: int = 8) -> dict:
    """Búsqueda unificada por contenido en Drive, Gmail y Calendar a la vez.

    Busca `texto` simultáneamente en los tres servicios de Google y devuelve
    un panorama común del despacho:
      - Drive: documentos cuyo CONTENIDO contiene el texto (fullText contains).
      - Gmail: correos que coincidan con la consulta (con De/Asunto/Fecha).
      - Calendar: eventos que coincidan (resumen, fecha y enlace).

    Args:
        texto: Término o frase a buscar. Es la base de "pregúntale a todo tu despacho".
        max_por_fuente: Máximo de resultados por cada fuente (Drive, Gmail, Calendar).

    Returns:
        dict con las listas 'archivos', 'correos' y 'eventos', más el recuento total.
    """
    try:
        archivos = []
        correos = []
        eventos = []

        # (a) Drive: búsqueda por contenido. Escapamos comillas simples para la query.
        texto_drive = texto.replace("'", "\\'")
        drive = _get_service()
        resp = drive.files().list(
            q="fullText contains '" + texto_drive + "' and trashed=false",
            fields="files(" + FILE_FIELDS + ")",
            pageSize=max_por_fuente,
        ).execute()
        for meta in resp.get("files", []):
            f = _file(meta)
            archivos.append({
                "nombre": f.get("name") or meta.get("name"),
                "tipo": f.get("mimeType") or meta.get("mimeType"),
                "enlace": f.get("webViewLink") or meta.get("webViewLink"),
            })

        # (b) Gmail: mensajes que coincidan y sus metadatos From/Subject/Date.
        gmail = _get_gmail()
        lista = gmail.users().messages().list(
            userId="me", q=texto, maxResults=max_por_fuente
        ).execute()
        for m in lista.get("messages", []):
            det = gmail.users().messages().get(
                userId="me", id=m["id"], format="metadata",
                metadataHeaders=["From", "Subject", "Date"],
            ).execute()
            cabeceras = {}
            for h in det.get("payload", {}).get("headers", []):
                cabeceras[h.get("name")] = h.get("value")
            correos.append({
                "id": m["id"],
                "de": cabeceras.get("From"),
                "asunto": cabeceras.get("Subject"),
                "fecha": cabeceras.get("Date"),
            })

        # (c) Calendar: eventos que coincidan con el texto.
        cal = _get_cal()
        ev = cal.events().list(
            calendarId="primary", q=texto, singleEvents=True,
            orderBy="startTime", maxResults=max_por_fuente,
        ).execute()
        for item in ev.get("items", []):
            inicio = item.get("start", {})
            fecha = inicio.get("dateTime") or inicio.get("date")
            eventos.append({
                "id": item.get("id"),
                "resumen": item.get("summary"),
                "fecha": fecha,
                "enlace": item.get("htmlLink"),
            })

        return {
            "consulta": texto,
            "archivos": archivos,
            "correos": correos,
            "eventos": eventos,
            "total": len(archivos) + len(correos) + len(eventos),
        }
    except Exception as e:
        return _err("buscar_en_todo", e)


# =========================================================================== #
# NUEVAS HERRAMIENTAS (lote sin permisos nuevos): Gmail filtros/firma,
# transcripciones de Meet y escaner de plazos en documentos.
# =========================================================================== #

@mcp.tool()
def gmail_crear_filtro(de: Optional[str] = None, para: Optional[str] = None,
                       asunto: Optional[str] = None, contiene: Optional[str] = None,
                       tiene_adjunto: bool = False, etiquetar: Optional[str] = None,
                       archivar: bool = False, marcar_leido: bool = False) -> dict:
    """Crea una REGLA/FILTRO automatico en Gmail que se aplica a los correos entrantes.

    Criterios (al menos uno): `de` (remitente), `para` (destinatario), `asunto`,
    `contiene` (consulta libre estilo Gmail), `tiene_adjunto`.
    Acciones (al menos una): `etiquetar` (nombre de etiqueta; se crea si no existe),
    `archivar` (saca de la bandeja de entrada), `marcar_leido`.
    Ejemplo: de='inspeccion@...', etiquetar='ITSS', archivar=True."""
    try:
        svc = _get_gmail()
        crit = {}
        if de: crit["from"] = de
        if para: crit["to"] = para
        if asunto: crit["subject"] = asunto
        if contiene: crit["query"] = contiene
        if tiene_adjunto: crit["hasAttachment"] = True
        if not crit:
            return {"error": "Indica al menos un criterio (de, para, asunto, contiene o tiene_adjunto)."}
        action = {}
        if etiquetar:
            labels = svc.users().labels().list(userId="me").execute().get("labels", [])
            lid = next((l["id"] for l in labels if l["name"].lower() == etiquetar.lower()), None)
            if not lid:
                lid = svc.users().labels().create(userId="me", body={
                    "name": etiquetar, "labelListVisibility": "labelShow",
                    "messageListVisibility": "show"}).execute()["id"]
            action["addLabelIds"] = [lid]
        rem = []
        if archivar: rem.append("INBOX")
        if marcar_leido: rem.append("UNREAD")
        if rem: action["removeLabelIds"] = rem
        if not action:
            return {"error": "Indica una accion (etiquetar, archivar o marcar_leido)."}
        r = svc.users().settings().filters().create(
            userId="me", body={"criteria": crit, "action": action}).execute()
        _audit("gmail_crear_filtro", {"id": r.get("id"), "criteria": crit})
        return {"ok": True, "id": r.get("id"), "criterios": crit, "accion": action}
    except Exception as e:
        return _err("gmail_crear_filtro", e)


@mcp.tool()
def gmail_listar_filtros() -> dict:
    """Lista los filtros/reglas automaticas configuradas en Gmail (id, criterios y accion)."""
    try:
        svc = _get_gmail()
        res = svc.users().settings().filters().list(userId="me").execute()
        out = []
        for f in res.get("filter", []):
            out.append({"id": f.get("id"), "criterios": f.get("criteria"),
                        "accion": f.get("action")})
        return {"filtros": out, "total": len(out)}
    except Exception as e:
        return _err("gmail_listar_filtros", e)


@mcp.tool()
def gmail_eliminar_filtro(filtro_id: str) -> dict:
    """Elimina un filtro/regla de Gmail por su id (ver gmail_listar_filtros)."""
    try:
        svc = _get_gmail()
        svc.users().settings().filters().delete(userId="me", id=filtro_id).execute()
        _audit("gmail_eliminar_filtro", {"id": filtro_id})
        return {"ok": True, "id": filtro_id}
    except Exception as e:
        return _err("gmail_eliminar_filtro", e)


@mcp.tool()
def gmail_definir_firma(texto: str, email: Optional[str] = None) -> dict:
    """Define la FIRMA del correo (HTML o texto) para tu direccion principal o la
    direccion `email` indicada. Sustituye la firma actual de esa direccion."""
    try:
        svc = _get_gmail()
        sendas = svc.users().settings().sendAs().list(userId="me").execute().get("sendAs", [])
        addr = email or next((s["sendAsEmail"] for s in sendas if s.get("isPrimary")), None)
        if not addr:
            return {"error": "No encuentro la direccion de envio principal."}
        svc.users().settings().sendAs().patch(
            userId="me", sendAsEmail=addr, body={"signature": texto}).execute()
        _audit("gmail_definir_firma", {"email": addr})
        return {"ok": True, "email": addr}
    except Exception as e:
        return _err("gmail_definir_firma", e)


@mcp.tool()
def gmail_ver_firma(email: Optional[str] = None) -> dict:
    """Muestra la firma configurada para tu direccion principal (o la indicada)."""
    try:
        svc = _get_gmail()
        sendas = svc.users().settings().sendAs().list(userId="me").execute().get("sendAs", [])
        target = None
        for s in sendas:
            if (email and s.get("sendAsEmail") == email) or (not email and s.get("isPrimary")):
                target = s; break
        if not target:
            return {"error": "No encuentro esa direccion de envio."}
        return {"email": target.get("sendAsEmail"), "firma": target.get("signature", "")}
    except Exception as e:
        return _err("gmail_ver_firma", e)


@mcp.tool()
def meet_transcripciones(dias: int = 30, contiene: str = "") -> dict:
    """Localiza en tu Drive las TRANSCRIPCIONES y notas que Google Meet deja tras las
    reuniones (Google las guarda como Documentos de Google, normalmente en la carpeta
    'Meet Recordings'). Devuelve los documentos candidatos de los ultimos `dias`; si se
    da `contiene`, filtra por su contenido. Luego puedes leer uno con drive_read_file
    y redactar el acta."""
    try:
        import datetime as _dt
        svc = _get_service()
        since = (_dt.datetime.utcnow() - _dt.timedelta(days=max(1, dias))).strftime("%Y-%m-%dT%H:%M:%SZ")
        base = "mimeType='application/vnd.google-apps.document' and trashed=false and modifiedTime > '" + since + "'"
        if contiene:
            c = contiene.replace("'", "\\'")
            q = base + " and fullText contains '" + c + "'"
        else:
            terms = ["Transcript", "Transcripción", "Notas de la reunión",
                     "Gemini", "Recording", "Grabación"]
            q = base + " and (" + " or ".join("name contains '" + t + "'" for t in terms) + ")"
        res = svc.files().list(q=q, orderBy="modifiedTime desc",
              fields="files(id,name,modifiedTime,webViewLink)", pageSize=25,
              includeItemsFromAllDrives=True, supportsAllDrives=True).execute()
        return {"transcripciones": res.get("files", []), "total": len(res.get("files", []))}
    except Exception as e:
        return _err("meet_transcripciones", e)


@mcp.tool()
def documento_escanear_plazos(file_id: str) -> dict:
    """Lee un documento de Drive (resolucion, notificacion, requerimiento...) y DETECTA
    fechas y plazos en su texto: fechas explicitas (dd/mm/aaaa o 'dd de mes de aaaa') y
    expresiones del tipo 'en el plazo de N dias (habiles/naturales)'. Devuelve lo hallado
    para que luego crees los eventos y tareas oportunos. No calcula por si mismo el
    vencimiento definitivo: sirve de apoyo, no sustituye la verificacion del profesional."""
    try:
        import re as _re
        r = drive_read_file(file_id)
        texto = r.get("content") or r.get("text") or ""
        if not texto:
            return {"error": "No pude extraer texto del documento.", "detalle": r}
        meses = "enero|febrero|marzo|abril|mayo|junio|julio|agosto|septiembre|octubre|noviembre|diciembre"
        fechas = []
        for m in _re.finditer(r"\b(\d{1,2})[/-](\d{1,2})[/-](\d{2,4})\b", texto):
            fechas.append(m.group(0))
        for m in _re.finditer(r"\b\d{1,2}\s+de\s+(?:" + meses + r")\s+de\s+\d{4}\b", texto, _re.IGNORECASE):
            fechas.append(m.group(0))
        plazos = []
        for m in _re.finditer(r"plazo\s+de\s+(\w+|\d+)\s+d[ií]as(?:\s+(h[aá]biles|naturales))?", texto, _re.IGNORECASE):
            plazos.append(m.group(0).strip())
        # dedupe conservando orden
        def uniq(xs):
            seen = set(); out = []
            for x in xs:
                k = x.lower()
                if k not in seen:
                    seen.add(k); out.append(x)
            return out
        return {"ok": True, "nombre": r.get("name"),
                "fechas_detectadas": uniq(fechas),
                "plazos_detectados": uniq(plazos),
                "aviso": "Deteccion automatica de apoyo; verifica cada plazo en el procedimiento concreto."}
    except Exception as e:
        return _err("documento_escanear_plazos", e)


# --------------------------------------------------------------------------- #
# Arranque
# --------------------------------------------------------------------------- #

if __name__ == "__main__":
    mcp.run(transport="streamable-http")
